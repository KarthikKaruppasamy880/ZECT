"""Present + Voice production gates — honest BLOCKED_EXTERNAL, never fake PASS."""

from __future__ import annotations

import io
import os
import uuid
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.infrastructure.database import SessionLocal
from app.models import ClonedVoice, User
from app.services.mentrix.presentation.final_pptx_inspector import inspect_pptx_bytes
from app.services.mentrix.presentation.quality_critic import critique_plan
from app.services.mentrix.presentation.renderer import render_plan_to_pptx, validate_generated_pptx
from app.services.mentrix.presentation.template_importer import UnsafePptxError
from app.services.mentrix.presentation.deck_catalog import import_pptx_bytes
from tests.fixes_and_phases.pptx_fixtures import make_master_pptx_bytes
from tests.fixes_and_phases.test_visual_parity import _visual_plan
from app.services.mentrix.presentation.asset_resolver import example_png_bytes, store_image
from app.services.mentrix.presentation.plan import validate_plan


def _overlapping_dump_pptx() -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "AI Agentic vs Graph, Loop, KV"
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(7), Inches(1.1))
    box.text_frame.text = "Agentic, graph, loop, KV cache, LLM fine-tuning: key AI terms."
    dump = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(6.5))
    dump.text_frame.text = (
        "AI Agentic vs Graph, Loop, KV Agentic, graph, loop, KV cache, LLM fine-tuning: key AI terms."
    )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_import_rejects_garbage_and_non_pptx_zip():
    with pytest.raises(UnsafePptxError):
        import_pptx_bytes(b"PK\x03\x04fake-pptx", filename="fake.pptx")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("readme.txt", "not a presentation")
    with pytest.raises(UnsafePptxError):
        import_pptx_bytes(buf.getvalue(), filename="empty.zip.pptx")


def test_blank_import_notes_idempotent_quality(authed_client):
    blank = authed_client.post("/api/mentrix/present/blank")
    assert blank.status_code == 200, blank.text
    path = blank.json()["path"]
    assert Path(path).is_file()

    decks = authed_client.get("/api/mentrix/present/decks")
    assert decks.status_code == 200
    names = {row.get("name") for row in decks.json().get("items") or []}
    assert Path(path).name in names

    parsed = authed_client.post("/api/mentrix/present/parse-pptx-path", json={"path": path})
    assert parsed.status_code == 200, parsed.text
    slides = parsed.json()["slides"]
    assert len(slides) >= 1
    notes_payload = []
    for row in slides:
        notes_payload.append(
            {
                "index": row.get("index", 0),
                "text": row.get("text") or "Untitled presentation",
                "notes": "Executive note: owners needed this week.",
            }
        )
    saved = authed_client.post("/api/mentrix/present/save-notes", json={"path": path, "slides": notes_payload})
    assert saved.status_code == 200, saved.text
    assert saved.json().get("ok") is True
    saved2 = authed_client.post("/api/mentrix/present/save-notes", json={"path": path, "slides": notes_payload})
    assert saved2.status_code == 200
    again = authed_client.post("/api/mentrix/present/parse-pptx-path", json={"path": path})
    assert again.status_code == 200
    blob = " ".join(str(s.get("notes") or "") for s in again.json()["slides"])
    assert "owners needed this week" in blob

    imported = authed_client.post(
        "/api/mentrix/present/import",
        files={
            "file": (
                f"reimport-{uuid.uuid4().hex[:8]}.pptx",
                Path(path).read_bytes(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert imported.status_code == 200, imported.text
    dest = imported.json()["path"]
    gate = authed_client.get("/api/mentrix/present/quality-gate", params={"path": dest})
    assert gate.status_code == 200, gate.text
    body = gate.json()
    assert body.get("ok") is True
    assert body.get("export_blocked") is not True, body
    first = authed_client.get("/api/mentrix/present/pptx", params={"path": dest})
    second = authed_client.get("/api/mentrix/present/pptx", params={"path": dest})
    assert first.status_code == 200, first.text
    assert second.status_code == 200
    assert len(first.content) == len(second.content)
    assert len(first.content) > 100


def test_import_http_rejects_invalid_pptx(authed_client):
    bad = authed_client.post(
        "/api/mentrix/present/import",
        files={"file": ("fake.pptx", b"PK\x03\x04fake-pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    assert bad.status_code == 400
    not_pptx = authed_client.post(
        "/api/mentrix/present/import",
        files={"file": ("notes.txt", b"hello", "text/plain")},
    )
    assert not_pptx.status_code == 400


def test_critical_quality_409_not_overridable(authed_client):
    from app.services.pptx_paths import default_pptx_save_dir

    dest = default_pptx_save_dir() / f"_prod_hard_block_{uuid.uuid4().hex[:8]}.pptx"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(_overlapping_dump_pptx())
    try:
        report = inspect_pptx_bytes(dest.read_bytes())
        assert report["status"] == "FAIL"
        gate = authed_client.get("/api/mentrix/present/quality-gate", params={"path": str(dest)})
        assert gate.status_code == 200, gate.text
        body = gate.json()
        assert body["export_blocked"] is True
        assert body["accept_warnings_allowed"] is False
        blocked = authed_client.get(
            "/api/mentrix/present/pptx",
            params={"path": str(dest), "accept_warnings": "true"},
        )
        assert blocked.status_code == 409
        detail = blocked.json().get("detail") or {}
        if isinstance(detail, dict):
            assert detail.get("error") == "export_blocked_critical_quality"
    finally:
        dest.unlink(missing_ok=True)


def test_zinnia_and_user_templates_http(authed_client, tmp_path, monkeypatch):
    monkeypatch.setenv("ZECT_PRESENT_TEMPLATE_ROOT", str(tmp_path))
    listed = authed_client.get("/api/mentrix/presentation/templates")
    assert listed.status_code == 200, listed.text
    body = listed.json()
    ids = {t["id"] for t in body.get("zinnia") or []}
    assert "zinnia-executive-v1" in ids
    assert "zinnia-delivery-v1" in ids
    uploaded = authed_client.post(
        "/api/mentrix/presentation/templates/upload",
        files={"file": ("mine.pptx", make_master_pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        data={"name": "User Prod Template", "scope": "USER"},
    )
    assert uploaded.status_code == 200, uploaded.text
    assert uploaded.json().get("ok") is True
    tid = uploaded.json()["template"]["id"]
    assert tid.startswith("user-")
    again = authed_client.get("/api/mentrix/presentation/templates")
    mine = {t["id"] for t in again.json().get("my_templates") or []}
    assert tid in mine


def test_critic_case_prompts_do_not_crash():
    cases = [
        ("executive", "Zinnia executive brief: Q3 delivery status, top risks, decisions needed."),
        ("roadmap", "Delivery roadmap: milestones, owners, and remaining work."),
        ("architecture", "Architecture overview: API fabric, control plane, and runtime."),
        ("metrics", "Metrics review: latency, error rate, and queue depth. Do not invent numbers."),
    ]
    for _label, prompt in cases:
        plan = {
            "slides": [
                {
                    "title": _label.title(),
                    "content_blocks": [{"kind": "bullet", "text": prompt}],
                    "blocks": [],
                    "notes_intent": prompt,
                }
            ]
        }
        report = critique_plan(plan, None, prompt=prompt)
        assert report.get("slides")
        assert "status" in report or "overlap_count" in report


def test_native_visual_render_inspector_not_presenton(tmp_path, monkeypatch):
    monkeypatch.setenv("ZECT_PRESENT_ASSET_ROOT", str(tmp_path / "assets"))
    meta = store_image(example_png_bytes(label="prod"), user_id="prod", filename="fig.png")
    plan = validate_plan(
        _visual_plan(image_asset=meta["asset_id"]),
        n_slides=4,
        template_id="",
        audience_id="executive",
    )
    data = render_plan_to_pptx(plan, user_id="prod")
    validate_generated_pptx(data, n_slides=4)
    report = inspect_pptx_bytes(data)
    assert report["slide_count"] == 4
    assert int(report.get("broken_rel_count") or 0) == 0
    from app.services.mentrix.presentation.document import inspect_pptx_visuals

    visuals = inspect_pptx_visuals(data)
    assert visuals["has_chart"] is True
    assert visuals["has_table"] is True
    assert visuals["has_image"] is True
    # Valid OOXML + visuals is not live Presenton Quality PASS.
    assert report.get("status") in {"PASS", "WARN", "FAIL"}


def test_presenton_status_honest(authed_client):
    st = authed_client.get("/api/mentrix/presenton/status")
    assert st.status_code == 200, st.text
    body = st.json()
    assert "lifecycle" in body
    assert "blocked_external" in body
    if body.get("blocked_external") or str(body.get("lifecycle") or "") in {
        "PROVIDER_UNAVAILABLE",
        "STARTING",
        "TEMPLATE_NOT_READY",
    }:
        # Honest unavailable/not-ready is not a generate PASS.
        return
    if not os.getenv("ZECT_LIVE_PRESENT"):
        pytest.skip("Presenton reports reachable; live Quality generate is opt-in ZECT_LIVE_PRESENT=1")


def test_voice_engine_status_honest(authed_client):
    st = authed_client.get("/api/mentrix/voice/engine-status")
    assert st.status_code == 200, st.text
    body = st.json()
    assert isinstance(body.get("online"), bool)
    if not body.get("online"):
        assert "Voicebox" in str(body.get("hint") or "") or "offline" in str(body.get("hint") or "").lower()


def test_voice_cross_user_http_denied(authed_client):
    vid = f"victim-voice-{uuid.uuid4().hex[:10]}"
    db = SessionLocal()
    try:
        victim = db.query(User).filter(User.email == "voice-victim@zect.local").first()
        if victim is None:
            victim = User(email="voice-victim@zect.local", name="Voice Victim", role="developer")
            db.add(victim)
            db.commit()
            db.refresh(victim)
        db.add(
            ClonedVoice(
                user_id=int(victim.id),
                voice_id=vid,
                name="Victim",
                provider="chatterbox",
                is_default=True,
                sample_path="",
                reference_text="hello",
            )
        )
        db.commit()
    finally:
        db.close()

    listed = authed_client.get("/api/mentrix/voice/voices")
    assert listed.status_code == 200, listed.text
    ids = {row.get("voice_id") for row in listed.json()}
    assert vid not in ids

    speak = authed_client.post("/api/mentrix/voice/speak", json={"text": "hello", "voice_id": vid})
    assert speak.status_code == 404

    delete = authed_client.delete(f"/api/mentrix/voice/voices/{vid}")
    assert delete.status_code == 404


def test_stock_voice_invalid_and_unconfigured(authed_client):
    bad = authed_client.post("/api/mentrix/voice/speak", json={"text": "hello", "stock_voice": "not-a-voice"})
    assert bad.status_code == 400
    from app.adapters.llm.openai_tts import openai_tts_available

    stock = authed_client.post("/api/mentrix/voice/speak", json={"text": "hello", "stock_voice": "nova"})
    if not openai_tts_available():
        assert stock.status_code == 503
        return
    if not os.getenv("ZECT_LIVE_VOICE_STOCK"):
        pytest.skip("OPENAI stock TTS configured; live stock speak is opt-in ZECT_LIVE_VOICE_STOCK=1")
    assert stock.status_code == 200
    assert len(stock.content) > 100


def test_powerpoint_com_opt_in_or_blocked_external():
    if not os.getenv("ZECT_LIVE_PPT_COM"):
        pytest.skip("BLOCKED_EXTERNAL: live Microsoft PowerPoint COM is opt-in ZECT_LIVE_PPT_COM=1")
    if os.name != "nt":
        pytest.skip("BLOCKED_EXTERNAL: PowerPoint COM is Windows-only")
    try:
        import win32com.client  # type: ignore
    except ImportError:
        pytest.skip("BLOCKED_EXTERNAL: pywin32 not installed")
    from app.services.pptx_paths import default_pptx_save_dir

    dest = default_pptx_save_dir() / f"_prod_com_{uuid.uuid4().hex[:8]}.pptx"
    dest.write_bytes(make_master_pptx_bytes())
    app = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(str(dest), WithWindow=False)
        assert int(pres.Slides.Count) >= 1
        pres.Close()
    finally:
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        dest.unlink(missing_ok=True)
