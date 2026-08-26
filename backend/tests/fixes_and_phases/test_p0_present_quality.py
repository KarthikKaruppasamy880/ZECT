"""P0 Present export quality: OOXML inspector, composition mutual exclusion, keep-cleanup."""

from __future__ import annotations

import inspect
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.infrastructure.database import SessionLocal
from app.models import Project
from app.services.fixture_isolation import keep_cleanup_projects
from app.services.mentrix.presentation.document_io import apply_document_to_pptx
from app.services.mentrix.presentation.final_pptx_inspector import (
    inspect_and_repair_pptx,
    inspect_pptx_bytes,
    strip_covering_dump_shapes,
)
from app.services.mentrix.presentation.layout_composer import compose_regions
from app.services.mentrix.presentation.native_provider import ZectNativePresentationProvider
from app.services.mentrix.presentation.quality_critic import critique_plan
from app.services.mentrix.presentation.quality_repair import repair_until_pass
from app.services.mentrix.presentation.renderer import render_plan_to_pptx


def _overlapping_dump_pptx() -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "AI Agentic vs Graph, Loop, KV"
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(7), Inches(1.1))
    box.text_frame.text = "Agentic, graph, loop, KV cache, LLM fine-tuning: key AI terms."
    dump = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(6.5))
    dump.text_frame.text = (
        "AI Agentic vs Graph, Loop, KV Agentic, graph, loop, KV cache, LLM fine-tuning: key AI terms."
    )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _fixture_pptx() -> Path | None:
    here = Path(__file__).resolve().parent
    for candidate in (
        here / "fixtures" / "zect-deck-overlap.pptx",
        here.parents[2] / "prompts" / "zect-deck.pptx",
    ):
        if candidate.is_file():
            return candidate
    return None


def test_inspector_flags_covering_dump_and_strip_repairs():
    data = _overlapping_dump_pptx()
    report = inspect_pptx_bytes(data)
    assert report["covering_dump_count"] >= 1 or report["overlap_count"] >= 1
    assert report["status"] == "FAIL"
    repaired, removed = strip_covering_dump_shapes(data)
    assert removed >= 1
    after = inspect_pptx_bytes(repaired)
    assert after["covering_dump_count"] == 0


def test_zect_deck_fixture_duplicate_overlap():
    path = _fixture_pptx()
    if path is None:
        return
    report = inspect_pptx_bytes(path.read_bytes())
    assert report["slide_count"] >= 1
    assert report["status"] == "FAIL"
    assert report["export_blocked"] is True
    assert report["covering_dump_count"] >= 1 or report["overlap_count"] >= 1
    assert report.get("hard_findings")
    repaired, _n = strip_covering_dump_shapes(path.read_bytes())
    after = inspect_pptx_bytes(repaired)
    assert after["covering_dump_count"] == 0


def test_critic_flags_duplicate_overlapping_text():
    shared = {"x": 400000, "y": 1200000, "cx": 8000000, "cy": 3000000}
    closing = "Thank you for your attention."
    plan = {
        "slides": [
            {
                "title": "Closing",
                "composed_regions": {
                    "title": {"x": 400000, "y": 200000, "cx": 8000000, "cy": 800000},
                    "body": shared,
                },
                "content_blocks": [],
                "blocks": [
                    {"kind": "text", "content": {"text": closing}, "geometry": shared},
                    {"kind": "text", "content": {"text": closing}, "geometry": shared},
                ],
            }
        ]
    }
    report = critique_plan(plan, None, prompt="closing slide")
    assert report["duplicate_semantic_count"] >= 1
    assert report["status"] == "FAIL"
    slide_findings = report["slides"][0]["findings"]
    assert "duplicate_semantic_content" in slide_findings


def test_zinnia_master_does_not_stack_title_textbox_on_object_placeholder():
    root = Path(__file__).resolve().parents[2]
    master = root.parent / ".zect/present-templates/masters/zinnia-executive-v1.pptx"
    if not master.is_file():
        return
    plan = {
        "slides": [
            {
                "title": "Q3 Delivery Overview",
                "content_blocks": [
                    {"kind": "bullet", "text": "Overview of Q3 objectives and deliverables."},
                    {"kind": "bullet", "text": "Summary of progress against targets."},
                ],
                "blocks": [
                    {
                        "kind": "text",
                        "content": {"text": "Overview of Q3 objectives and deliverables."},
                        "validation": {"ok": True, "errors": []},
                    }
                ],
            }
        ]
    }
    data = render_plan_to_pptx(plan, template_path=master)
    report = inspect_pptx_bytes(data)
    findings = [f for slide in report["slides"] for f in slide.get("findings") or []]
    assert "placeholder_and_generated" not in findings
    prs = Presentation(io.BytesIO(data))
    slide = prs.slides[0]
    titled = []
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        text = (sh.text_frame.text or "").strip()
        if not text:
            continue
        titled.append((sh.name, int(sh.left), int(sh.top), int(sh.width), int(sh.height), text.splitlines()[0][:80]))
    for i, a in enumerate(titled):
        for b in titled[i + 1 :]:
            ax2, ay2 = a[1] + a[3], a[2] + a[4]
            bx2, by2 = b[1] + b[3], b[2] + b[4]
            overlap_x = min(ax2, bx2) - max(a[1], b[1])
            overlap_y = min(ay2, by2) - max(a[2], b[2])
            assert overlap_x <= 40000 or overlap_y <= 40000, (a[0], b[0], a[5], b[5])


def test_placeholder_and_generated_are_mutually_exclusive():
    plan = {
        "slides": [
            {
                "title": "Board opening",
                "content_blocks": [{"kind": "bullet", "text": "Q3 delivery is on track"}],
                "blocks": [
                    {
                        "kind": "text",
                        "content": {"text": "Q3 delivery is on track"},
                        "validation": {"ok": True, "errors": []},
                    }
                ],
            }
        ]
    }
    data = render_plan_to_pptx(plan)
    report = inspect_pptx_bytes(data)
    findings = [f for slide in report["slides"] for f in slide.get("findings") or []]
    assert "placeholder_and_generated" not in findings
    assert report["covering_dump_count"] == 0


def test_inspector_flags_layout_inherited_placeholder_overlap():
    """Slide XML often omits placeholder xfrms; export gate must still see COM-visible overlap."""
    master = Path(__file__).resolve().parents[2].parent / ".zect/present-templates/masters/zinnia-executive-v1.pptx"
    if not master.is_file():
        return
    prs = Presentation(str(master))
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass
        sld_id_lst.remove(sld_id)
    layout = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    obj = None
    for shape in slide.shapes:
        try:
            if int(shape.placeholder_format.type) == 7:
                obj = shape
                break
        except Exception:
            continue
    if obj is None:
        return
    obj.text_frame.text = "Overview of Q3 objectives and deliverables."
    box = slide.shapes.add_textbox(obj.left, obj.top, obj.width, obj.height)
    box.text_frame.text = "Q3 Delivery Overview"
    buf = io.BytesIO()
    prs.save(buf)
    report = inspect_pptx_bytes(buf.getvalue())
    assert report["status"] == "FAIL"
    assert "text_shape_collision" in report["hard_findings"]
    assert report["export_blocked"] is True


def test_compose_splits_body_and_visual():
    definition = {
        "slide_size": {"cx": 9144000, "cy": 5143500},
        "layouts": [
            {
                "name": "Title and Content",
                "placeholders": [
                    {"type": "TITLE", "geometry": {"x": 457200, "y": 200000, "cx": 8229600, "cy": 700000}},
                    {"type": "BODY", "geometry": {"x": 457200, "y": 1100000, "cx": 8229600, "cy": 3600000}},
                ],
            }
        ],
    }
    regions = compose_regions(definition, definition["layouts"][0], split_visual=True)
    assert regions["visual"]["y"] != regions["body"]["y"] or regions["visual"]["cy"] != regions["body"]["cy"]


def test_renderer_does_not_add_covering_dump_on_top_of_placeholders():
    plan = {
        "slides": [
            {
                "title": "Board opening",
                "content_blocks": [{"kind": "bullet", "text": "Q3 delivery is on track"}],
                "blocks": [
                    {
                        "kind": "text",
                        "content": {"text": "Q3 delivery is on track"},
                        "validation": {"ok": True, "errors": []},
                    }
                ],
            }
        ]
    }
    data = render_plan_to_pptx(plan)
    report = inspect_pptx_bytes(data)
    assert report["covering_dump_count"] == 0


def test_editor_save_export_idempotent_shape_count(tmp_path):
    data = render_plan_to_pptx(
        {
            "slides": [
                {
                    "title": "Status",
                    "content_blocks": [{"kind": "bullet", "text": "Twelve epics closed"}],
                    "blocks": [],
                }
            ]
        }
    )
    dest = tmp_path / "roundtrip.pptx"
    dest.write_bytes(data)
    before = inspect_pptx_bytes(data)["slides"][0]["shape_count"]
    apply_document_to_pptx(
        dest,
        [{"index": 0, "text": "Status\nTwelve epics closed", "notes": "say this"}],
    )
    after = inspect_pptx_bytes(dest.read_bytes())
    assert after["slides"][0]["shape_count"] <= before + 1
    assert after["covering_dump_count"] == 0
    apply_document_to_pptx(
        dest,
        [{"index": 0, "text": "Status\nTwelve epics closed", "notes": "say this"}],
    )
    again = inspect_pptx_bytes(dest.read_bytes())
    assert again["slides"][0]["shape_count"] == after["slides"][0]["shape_count"]


def test_inspect_and_repair_clears_legacy_overlap_fixture():
    path = _fixture_pptx()
    if path is None:
        return
    before = inspect_pptx_bytes(path.read_bytes())
    fixed, report = inspect_and_repair_pptx(path.read_bytes())
    after = inspect_pptx_bytes(fixed)
    assert before["status"] == "FAIL"
    assert after["status"] == "PASS"
    assert after["overlap_count"] == 0
    assert not after.get("hard_findings")


def test_semantic_duplicate_skips_short_label_in_body():
    from app.services.mentrix.presentation.final_pptx_inspector import _semantic_duplicate_text

    assert not _semantic_duplicate_text("repeat", "repeats adapt from earlier cycles")
    assert _semantic_duplicate_text("same title", "same title")


def test_split_title_spills_to_key_message():
    from app.services.mentrix.presentation.content_capacity import split_title_for_regions

    long_title = "Difference between AI Agentic and the Graph, loop and KV cache with LLM fine tuning"
    headline, subtitle = split_title_for_regions(long_title, "", {"max_title_chars": 48, "max_title_lines": 2})
    assert len(headline) < len(long_title)
    assert len(headline) <= 48
    assert headline in long_title or long_title.startswith(headline.split()[0])


def test_render_populates_subtitle_and_clears_date_sample():
    master = Path(__file__).resolve().parents[2].parent / ".zect/present-templates/masters/zinnia-executive-v1.pptx"
    if not master.is_file():
        return
    plan = {
        "slides": [
            {
                "title": "AI Agentic Systems Overview",
                "key_message": "How agentic AI differs from graph and loop patterns.",
                "purpose": "opening",
                "layout_intent": "title",
                "content_blocks": [],
                "blocks": [],
            }
        ]
    }
    data = render_plan_to_pptx(plan, template_path=master)
    assert b"Date Here" not in data
    report = inspect_pptx_bytes(data)
    findings = [f for slide in report["slides"] for f in slide.get("findings") or []]
    assert "placeholder_and_generated" not in findings


def test_fast_and_quality_both_call_inspector_and_critic():
    src = inspect.getsource(ZectNativePresentationProvider.generate)
    assert "repair_until_pass" in src
    assert "inspect_and_repair_pptx" in src
    data = render_plan_to_pptx(
        {"slides": [{"title": "A", "content_blocks": [{"kind": "bullet", "text": "One"}], "blocks": []}]}
    )
    _fixed, report = inspect_and_repair_pptx(data)
    assert "status" in report


def test_degraded_fast_does_not_override_layout_fail():
    plan = {
        "slides": [
            {
                "title": "Overlap",
                "composed_regions": {
                    "title": {"x": 400000, "y": 200000, "cx": 8000000, "cy": 800000},
                    "body": {"x": 400000, "y": 400000, "cx": 8000000, "cy": 3000000},
                },
                "content_blocks": [{"kind": "bullet", "text": "Body"}],
                "blocks": [],
            }
        ]
    }
    _plan, report = repair_until_pass(plan, None, prompt="status", degraded=True)
    layout_hard = bool(report.get("overlap_count") or report.get("out_of_bounds_count"))
    if layout_hard:
        assert report["status"] == "FAIL"
        assert not report.get("degraded_override")
        assert report.get("status") in {"FAIL", "NEEDS_REVIEW"}


def test_export_blocked_on_quality_failed():
    data = _overlapping_dump_pptx()
    report = inspect_pptx_bytes(data)
    assert report["status"] == "FAIL"
    from app.services.mentrix.presentation.deck_catalog import quality_gate_for_path
    from app.services.pptx_paths import default_pptx_save_dir

    dest = default_pptx_save_dir() / "_p0_quality_fail.pptx"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(data)
    gate = quality_gate_for_path(str(dest))
    assert gate["export_blocked"] is True
    assert gate["quality_passed"] is False
    assert gate["accept_warnings_allowed"] is False
    assert gate.get("hard_findings")
    dest.unlink(missing_ok=True)


def test_accept_warnings_cannot_override_critical_export_block(authed_client, tmp_path):
    from app.services.pptx_paths import default_pptx_save_dir

    dest = default_pptx_save_dir() / "_p0_hard_block.pptx"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(_overlapping_dump_pptx())
    gate = authed_client.get(f"/api/mentrix/present/quality-gate?path={dest}")
    assert gate.status_code == 200, gate.text
    body = gate.json()
    assert body["export_blocked"] is True
    assert body["accept_warnings_allowed"] is False
    blocked = authed_client.get(
        f"/api/mentrix/present/pptx?path={dest}&accept_warnings=true"
    )
    assert blocked.status_code == 409
    detail = blocked.json().get("detail") or {}
    if isinstance(detail, dict):
        assert detail.get("error") == "export_blocked_critical_quality"
    dest.unlink(missing_ok=True)


def test_compose_dedupes_duplicate_closing_text():
    from app.services.mentrix.presentation.layout_composer import compose_plan

    shared = {"x": 400000, "y": 1200000, "cx": 8000000, "cy": 3000000}
    closing = "Thank you for your attention."
    plan = {
        "slides": [
            {
                "title": "Closing",
                "content_blocks": [{"kind": "bullet", "text": closing}],
                "blocks": [
                    {"kind": "text", "content": {"text": closing}, "geometry": shared},
                    {"kind": "text", "content": {"text": closing}, "geometry": shared},
                ],
            }
        ]
    }
    compose_plan(plan, None)
    texts = [
        str((b.get("content") or {}).get("text") or "")
        for b in list(plan["slides"][0].get("blocks") or [])
        if str(b.get("kind") or "") in {"text", "bullet", "body"}
    ]
    assert texts.count(closing) <= 1
    assert int(plan["slides"][0].get("dedupe_removed") or 0) >= 1


def test_content_budget_trims_long_title():
    from app.services.mentrix.presentation.content_capacity import apply_content_budget

    regions = {
        "title": {"x": 400000, "y": 200000, "cx": 2000000, "cy": 400000},
        "body": {"x": 400000, "y": 900000, "cx": 8000000, "cy": 3000000},
    }
    slide = {"title": "X" * 200, "content_blocks": [{"text": "Y" * 300}]}
    cap = apply_content_budget(slide, regions)
    assert len(slide["title"]) <= cap["max_title_chars"]
    assert len(slide["content_blocks"][0]["text"]) <= cap["max_bullet_chars"]


def test_repair_dedupes_duplicate_semantic_fail():
    shared = {"x": 400000, "y": 1200000, "cx": 8000000, "cy": 3000000}
    closing = "Thank you for your attention."
    plan = {
        "slides": [
            {
                "title": "Closing",
                "composed_regions": {
                    "title": {"x": 400000, "y": 200000, "cx": 8000000, "cy": 800000},
                    "body": shared,
                },
                "content_blocks": [],
                "blocks": [
                    {"kind": "text", "content": {"text": closing}, "geometry": shared},
                    {"kind": "text", "content": {"text": closing}, "geometry": shared},
                ],
            }
        ]
    }
    _plan, report = repair_until_pass(plan, None, prompt="closing")
    assert int(report.get("duplicate_semantic_count") or 0) == 0 or report.get("status") == "PASS"


def test_kv_cache_grounding_wording():
    plan = {
        "slides": [
            {
                "title": "KV cache",
                "content_blocks": [
                    {"kind": "bullet", "text": "KV cache reduces memory use in the transformer."},
                ],
                "blocks": [],
            }
        ]
    }
    report = critique_plan(plan, None, prompt="Explain KV cache vs recomputation")
    findings = [f for slide in report["slides"] for f in slide.get("findings") or []]
    assert "kv_cache_memory_oversimplified" in findings
    ok = {
        "slides": [
            {
                "title": "KV cache",
                "content_blocks": [
                    {
                        "kind": "bullet",
                        "text": "KV cache avoids recomputing attention keys; it trades memory for compute.",
                    }
                ],
                "blocks": [],
            }
        ]
    }
    ok_report = critique_plan(ok, None, prompt="Explain KV cache vs recomputation")
    ok_findings = [f for slide in ok_report["slides"] for f in slide.get("findings") or []]
    assert "kv_cache_memory_oversimplified" not in ok_findings


def test_keep_cleanup_never_deletes_keep_ids():
    db = SessionLocal()
    created = []
    try:
        keep = Project(name="P0 Keep Probe", provenance="user", test_run_id="")
        drop = Project(name="P0 Drop Probe", provenance="user", test_run_id="")
        db.add_all([keep, drop])
        db.commit()
        db.refresh(keep)
        db.refresh(drop)
        created = [int(keep.id), int(drop.id)]
        keep_ids = [int(p.id) for p in db.query(Project).all() if int(p.id) != int(drop.id)]
        dry = keep_cleanup_projects(db, keep_ids, dry_run=True)
        assert dry["ok"] is True
        ids = {row["id"] for row in dry["would_delete"]}
        assert drop.id in ids
        assert keep.id not in ids
        live = keep_cleanup_projects(db, keep_ids, dry_run=False)
        assert drop.id in live["deleted_ids"]
        assert keep.id not in live["deleted_ids"]
        db.expire_all()
        assert db.query(Project).filter(Project.id == keep.id).first() is not None
        assert db.query(Project).filter(Project.id == drop.id).first() is None
    finally:
        for pid in created:
            row = db.query(Project).filter(Project.id == pid).first()
            if row is not None:
                db.delete(row)
        db.commit()
        db.close()


def test_keep_cleanup_refuses_empty_keep_ids(authed_client):
    empty = authed_client.post("/api/projects/fixtures/keep-cleanup", json={"keep_ids": [], "dry_run": True})
    assert empty.status_code == 400
    keep = authed_client.post(
        "/api/projects",
        json={"name": "Keep Cleanup Probe", "description": "tmp", "provenance": "test", "test_run_id": "p0-keep"},
    )
    assert keep.status_code == 201
    kid = keep.json()["id"]
    dry = authed_client.post(
        "/api/projects/fixtures/keep-cleanup",
        json={"keep_ids": [kid], "dry_run": True},
    )
    assert dry.status_code == 200
    body = dry.json()
    assert body["ok"] is True
    assert body["dry_run"] is True
    assert kid not in {row["id"] for row in body["would_delete"]}
