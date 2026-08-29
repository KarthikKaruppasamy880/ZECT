#!/usr/bin/env python3
"""Golden V2 generation proof — requested 3 slides, Zinnia template, agentic topic."""

from __future__ import annotations

import json
import os
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

GOLDEN_PROMPT = (
    "AI Agentic architectures: graph-based agents, tool loops, and KV-cache efficiency "
    "for enterprise inference. Executive briefing for platform engineering."
)
REQUESTED = 3
TEMPLATE = "zinnia-executive-v1"


def main() -> int:
    os.environ.setdefault("ZECT_PRESENTATION_PROVIDER", "zect_native")
    from app.services.mentrix.presentation.native_provider import ZectNativePresentationProvider
    from app.services.mentrix.presentation.provider import PresentationGenerateRequest
    from app.services.pptx_parse import parse_pptx_bytes
    from pptx import Presentation

    provider = ZectNativePresentationProvider()
    out = provider.generate(
        PresentationGenerateRequest(
            content=GOLDEN_PROMPT,
            n_slides=REQUESTED,
            template=TEMPLATE,
            ui_template_choice=TEMPLATE,
            filename="golden-v2-agentic-deck.pptx",
            user_id="golden-proof",
            fast_basic=True,
        )
    )
    artifact_dir = ROOT / "artifacts" / "present-golden-v2"
    artifact_dir.mkdir(parents=True, exist_ok=True)
    report = {
        "ok": bool(out.get("ok")),
        "requested_slide_count": REQUESTED,
        "plan_n_slides": out.get("n_slides"),
        "generation_job_id": out.get("generation_job_id"),
        "slide_count_trace": out.get("slide_count_trace"),
        "final_quality_status": out.get("final_quality_status"),
        "export_blocked": out.get("export_blocked"),
        "path": out.get("path"),
    }
    if out.get("path"):
        pptx = Path(str(out["path"]))
        data = pptx.read_bytes()
        parsed = parse_pptx_bytes(data)
        prs = Presentation(str(pptx))
        report["pptx_slide_count"] = len(prs.slides)
        report["parse_slide_count"] = len(parsed)
        report["slide_count_match"] = len(prs.slides) == REQUESTED
    report["acceptance"] = bool(
        out.get("ok")
        and report.get("slide_count_match")
        and int(out.get("n_slides") or 0) == REQUESTED
    )
    out_path = artifact_dir / "golden_v2_report.json"
    out_path.write_text(json.dumps(report, indent=2), encoding="utf-8")
    print(json.dumps(report, indent=2))
    return 0 if report["acceptance"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
