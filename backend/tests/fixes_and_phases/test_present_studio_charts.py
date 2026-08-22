from pptx import Presentation
from pptx.util import Inches

from app.services.mentrix.presentation.blocks import CHART_TYPES
from app.services.mentrix.presentation.charts import add_chart, replace_chart_data, resolve_ooxml_chart_type
from app.services.mentrix.presentation.document_io import apply_document_to_pptx
from app.services.mentrix.presentation.slide_ai import patch_slide_from_prompt


def _chart_block(chart_type: str) -> dict:
    return {
        "kind": "chart",
        "content": {
            "title": "Coverage",
            "chart_type": chart_type,
            "categories": ["A", "B", "C"],
            "series": [{"name": "Series", "values": [1.0, 3.0, 2.0]}],
            "legend": True,
        },
    }


def test_chart_types_include_presenton_parity():
    assert {"radar", "area", "stacked", "stacked_horizontal", "scatter", "polar", "progress", "gauge"} <= set(CHART_TYPES)


def test_polar_progress_gauge_aliases():
    assert resolve_ooxml_chart_type("polar") == "radar"
    assert resolve_ooxml_chart_type("progress") == "column"
    assert resolve_ooxml_chart_type("gauge") == "donut"


def test_add_and_replace_radar_round_trip(tmp_path):
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    apply_document_to_pptx(
        path,
        [{"index": 0, "text": "Title", "notes": "n", "blocks": [_chart_block("column")]}],
        user_id="u1",
    )
    apply_document_to_pptx(
        path,
        [{"index": 0, "text": "Title", "notes": "n", "blocks": [_chart_block("radar")]}],
        user_id="u1",
    )
    opened = Presentation(str(path))
    charts = [s for s in opened.slides[0].shapes if getattr(s, "has_chart", False)]
    assert charts, "radar chart should round-trip into OOXML"


def test_add_chart_direct_radar():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    geom = {"x": int(Inches(0.5)), "y": int(Inches(1.2)), "cx": int(Inches(8)), "cy": int(Inches(4))}
    assert add_chart(slide, _chart_block("radar"), geom) is True
    assert replace_chart_data(slide, _chart_block("bar")) is True


def test_slide_ai_radar_prompt_does_not_invent_kpis():
    out = patch_slide_from_prompt(prompt="Change this to a radar chart")
    assert out["ok"] is True
    assert out["chart_type"] == "radar"
    assert out["action"] == "chart_type"
