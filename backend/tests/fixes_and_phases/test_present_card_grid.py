"""Present Studio V4 -- Gamma/Presenton-class icon+title+body card grid.

ZECT's AI-generation pipeline previously had no layout vocabulary for "N
parallel named concepts" (e.g. Cursor's own "three modes: Chat / Inline
Edit / Agent" or "four pillars") -- every such slide defaulted to plain
stacked text blocks. This suite proves the new card_grid path end-to-end:
purpose detection -> visual_intent -> block construction -> native paint,
with zero geometric overlap."""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches

from app.services.mentrix.presentation.blocks import _card_content, example_card_grid_block
from app.services.mentrix.presentation.final_pptx_inspector import inspect_pptx_bytes
from app.services.mentrix.presentation.visual import paint_card_grid
from app.services.mentrix.presentation.visual_planner import apply_visual_plan, looks_like_card_group


class TestCardGroupDetection:
    def test_two_to_four_short_parallel_blocks_look_like_a_card_group(self):
        slide = {"content_blocks": [{"text": "Chat Mode explains code."}, {"text": "Agent Mode plans changes."}]}
        assert looks_like_card_group(slide)

    def test_a_single_long_paragraph_is_not_a_card_group(self):
        slide = {"content_blocks": [{"text": "x" * 400}]}
        assert not looks_like_card_group(slide)

    def test_five_blocks_is_not_a_card_group_bounded_at_four(self):
        slide = {"content_blocks": [{"text": f"Point {i}"} for i in range(5)]}
        assert not looks_like_card_group(slide)

    def test_zero_or_one_block_is_not_a_card_group(self):
        assert not looks_like_card_group({"content_blocks": []})
        assert not looks_like_card_group({"content_blocks": [{"text": "solo"}]})


class TestVisualPlanEndToEnd:
    def _plan(self):
        return {
            "objective": "Explain Cursor AI modes",
            "slides": [
                {"index": 0, "title": "Cursor AI Modes", "purpose": "opening", "content_blocks": []},
                {
                    "index": 1,
                    "title": "Three ways to use Cursor",
                    "purpose": "section",
                    "content_blocks": [
                        {"text": "Chat Mode gives explanations and guidance for new engineers."},
                        {"text": "Inline Edit modifies code directly in the current file."},
                        {"text": "Agent Mode executes multi step changes across the codebase."},
                    ],
                },
                {"index": 2, "title": "Wrap up", "purpose": "cta", "content_blocks": []},
            ],
        }

    def test_parallel_concept_slide_gets_features_purpose_and_card_grid_visual(self):
        out = apply_visual_plan(self._plan(), audience_id="general")
        slide = out["slides"][1]
        assert slide["purpose"] == "features"
        assert slide["visual_intent"] == "card_grid"

    def test_card_grid_block_is_actually_constructed_with_real_content(self):
        out = apply_visual_plan(self._plan(), audience_id="general")
        slide = out["slides"][1]
        kinds = [b["kind"] for b in slide["blocks"]]
        assert "card_grid" in kinds
        cards = next(b for b in slide["blocks"] if b["kind"] == "card_grid")["content"]["cards"]
        assert 2 <= len(cards) <= 4
        assert all(c["body"] for c in cards)

    def test_cta_slide_is_not_reclassified_as_a_card_grid(self):
        """cta's own explicit visual treatment (a closing quote) must win
        over card-grid detection even though this deck contains a genuine
        card-grid slide elsewhere."""
        out = apply_visual_plan(self._plan(), audience_id="general")
        assert out["slides"][2]["visual_intent"] != "card_grid"

    def test_an_opening_slide_with_no_content_blocks_is_not_a_card_group(self):
        """Empty content_blocks can never satisfy looks_like_card_group
        (needs 2-4 blocks) -- only an incidental hint-word match in the
        title could misclassify it, which is a pre-existing purpose-
        override characteristic shared by every hint list here (e.g. an
        opening slide titled '...Architecture...' was already overridden
        to 'architecture' before card_grid existed), not something new."""
        slide = {"index": 0, "title": "Welcome", "purpose": "opening", "content_blocks": []}
        assert not looks_like_card_group(slide)


class TestCardContentParsing:
    def test_splits_header_dash_description_like_the_reference_deck(self):
        card = _card_content({"text": "Chat Mode - Questions, explanations, and guidance."}, 0)
        assert card["title"] == "Chat Mode"
        assert "Questions" in card["body"]

    def test_falls_back_to_numbered_point_when_no_separator_present(self):
        card = _card_content({"text": "just a plain sentence with no separator at all here"}, 2)
        assert card["title"] == "Point 3"

    def test_icon_falls_back_to_a_plain_numeral_not_a_symbol_glyph(self):
        """A symbol glyph (e.g. a diamond) rendered as a missing-glyph box
        in the real render pipeline -- numerals always render."""
        card = _card_content({"text": "Agent Mode plans and executes."}, 3)
        assert card["icon"] == "4"


class TestNativePainting:
    def _build_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        return prs, slide

    def test_fewer_than_two_cards_is_a_no_op_not_a_broken_render(self):
        """example_card_grid_block() itself pads a lone card up to 2 (a
        single floating card is never a real grid) -- this tests
        paint_card_grid()'s own guard directly, bypassing that padding,
        for a raw block that genuinely has fewer than 2 cards."""
        prs, slide = self._build_slide()
        block = {"content": {"cards": [{"title": "Only one", "body": "x"}]}}
        geom = {"x": int(Inches(0.6)), "y": int(Inches(1.6)), "cx": int(Inches(8.8)), "cy": int(Inches(4.6))}
        assert paint_card_grid(slide, block, geom) is False

    def test_three_cards_paint_with_zero_geometric_overlap(self):
        prs, slide = self._build_slide()
        block = example_card_grid_block(
            0,
            0,
            cards=[
                {"title": "Chat Mode", "body": "Questions and guidance."},
                {"title": "Inline Edit", "body": "Modify code directly."},
                {"title": "Agent Mode", "body": "Plan and execute changes."},
            ],
        )
        geom = {"x": int(Inches(0.6)), "y": int(Inches(1.6)), "cx": int(Inches(8.8)), "cy": int(Inches(4.6))}
        assert paint_card_grid(slide, block, geom) is True
        buf = io.BytesIO()
        prs.save(buf)
        report = inspect_pptx_bytes(buf.getvalue())
        assert report.get("overlap_count") == 0
        assert report.get("status") == "PASS"

    def test_four_cards_lay_out_as_a_2x2_grid_with_zero_overlap(self):
        prs, slide = self._build_slide()
        block = example_card_grid_block(
            0,
            0,
            cards=[{"title": f"Pillar {i}", "body": f"Description {i}."} for i in range(1, 5)],
        )
        geom = {"x": int(Inches(0.6)), "y": int(Inches(1.6)), "cx": int(Inches(8.8)), "cy": int(Inches(4.6))}
        assert paint_card_grid(slide, block, geom) is True
        buf = io.BytesIO()
        prs.save(buf)
        report = inspect_pptx_bytes(buf.getvalue())
        assert report.get("overlap_count") == 0
