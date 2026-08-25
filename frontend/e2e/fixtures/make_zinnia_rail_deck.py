"""Clone real Zinnia master and pad to >=22 slides for rail scroll acceptance."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

BACKEND = Path(__file__).resolve().parents[3] / "backend"
sys.path.insert(0, str(BACKEND))

TARGET_SLIDES = 22


def main() -> int:
    dest = Path(sys.argv[1]).resolve()
    from app.services.mentrix.presentation.deck_catalog import instantiate_from_template

    src = instantiate_from_template("zinnia-executive-v1")
    prs = Presentation(str(src))
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    while len(prs.slides) < TARGET_SLIDES:
        prs.slides.add_slide(blank)
    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    print(f"{dest} slides={len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
