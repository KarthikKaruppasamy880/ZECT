"""Native PPTX renderer — python-pptx (MIT). Does not vendor Presenton."""

from __future__ import annotations

import io
import re
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu, Inches

from app.services.mentrix.presentation.template_importer import UnsafePptxError, inspect_pptx_archive
from app.services.pptx_parse import parse_pptx_bytes

_FALLBACK_IDS = frozenset({"modern", "general", "standard", "swift", ""})


def _safe_filename(name: str) -> str:
    stem = re.sub(r"[^a-zA-Z0-9._-]+", "_", (name or "zect-deck").strip())[:80] or "zect-deck"
    if not stem.lower().endswith(".pptx"):
        stem += ".pptx"
    return stem


def _clear_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 — python-pptx has no public delete
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass
        sld_id_lst.remove(sld_id)


def _definition_layout_needles(definition: dict[str, Any] | None, intent: str) -> tuple[str, ...]:
    """Prefer TemplateDefinition layout names that match the slide intent."""
    names = [str(row.get("name") or "").strip().lower() for row in list((definition or {}).get("layouts") or [])]
    names = [n for n in names if n]
    want = (intent or "title_body").lower()
    prefer = {
        "title": ("title slide", "title"),
        "title_body": ("title and content", "title and body", "subtitle + 1", "1 column", "header", "content"),
        "two_column": ("two content", "comparison", "two column"),
        "section": ("section header", "section"),
        "closing": ("title slide", "blank"),
        "text_image": ("picture with caption", "title and content", "picture"),
        "full_image": ("blank", "picture"),
        "chart_commentary": ("title and content", "content"),
        "table": ("title and content", "content"),
        "comparison": ("comparison", "two content"),
        "metrics": ("title and content", "blank"),
        "quote": ("quote", "blank", "title slide"),
        "diagram": ("title and content", "blank"),
    }.get(want, ("content",))
    matched = tuple(n for n in names if any(p in n for p in prefer))
    return matched or prefer


def _pick_layout(
    prs: Presentation,
    intent: str,
    definition: dict[str, Any] | None = None,
    preferred_name: str = "",
):
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise UnsafePptxError("template_has_no_layouts")
    want_name = (preferred_name or "").strip().lower()
    named = [(sl, (sl.name or "").lower()) for sl in layouts]
    if want_name:
        for sl, name in named:
            if name == want_name:
                return sl
        for sl, name in named:
            if want_name in name or (name and name in want_name):
                return sl
    want = (intent or "title_body").lower()
    aliases = {
        "title": ("title slide", "title", "blank"),
        "title_body": ("subtitle + 1", "title + 1", "1 column", "header", "title and content", "title and body", "content", "text"),
        "two_column": ("two content", "comparison", "two column", "2 column", "subtitle + 2"),
        "section": ("section header", "section"),
        "closing": ("blank", "title slide", "end"),
        "text_image": ("title and content", "picture with caption", "content"),
        "full_image": ("blank", "picture", "title and content"),
        "chart_commentary": ("subtitle + 1", "1 column", "title and content", "content", "text"),
        "table": ("subtitle + 1", "1 column", "title and content", "content"),
        "comparison": ("comparison", "two content", "two column", "2 column"),
        "metrics": ("subtitle + 1", "title and content", "content", "blank"),
        "quote": ("quote", "blank", "title slide"),
        "diagram": ("blank", "title and content", "content"),
    }
    needles = _definition_layout_needles(definition, want) + aliases.get(want, ("content",))
    for needle in needles:
        for sl, name in named:
            if needle and needle in name:
                return sl
    for sl in layouts:
        try:
            kinds = [ph.placeholder_format.type for ph in sl.placeholders]
        except ValueError:
            kinds = []
        if PP_PLACEHOLDER.BODY in kinds or PP_PLACEHOLDER.OBJECT in kinds:
            return sl
    return layouts[0]


def _geom_box(regions: dict[str, Any] | None, key: str, fallback: tuple[float, float, float, float]):
    geom = (regions or {}).get(key) if isinstance(regions, dict) else None
    if isinstance(geom, dict) and int(geom.get("cx") or 0) > 0:
        return Emu(int(geom["x"])), Emu(int(geom["y"])), Emu(int(geom["cx"])), Emu(int(geom["cy"]))
    return Inches(fallback[0]), Inches(fallback[1]), Inches(fallback[2]), Inches(fallback[3])


def _is_placeholder(shape) -> bool:
    try:
        _ = shape.placeholder_format.type
        return True
    except ValueError:
        return False


_DATE_PATTERNS = (
    "date here",
    "click to add",
    "subtitle",
    "your title here",
    "presentation title",
)


def _is_sample_layout_text(text: str) -> bool:
    norm = _norm_line(text)
    if not norm:
        return False
    return any(p in norm for p in _DATE_PATTERNS)


def _clear_placeholder_sample_text(slide) -> None:
    """Wipe leftover master/layout sample text before a single canonical write."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not _is_placeholder(shape):
            continue
        try:
            shape.text_frame.clear()
        except Exception:
            try:
                shape.text_frame.text = ""
            except Exception:
                pass


def _clear_layout_sample_text(slide) -> None:
    """Clear non-placeholder layout sample strings (e.g. Date Here on title slides)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if _is_placeholder(shape):
            continue
        try:
            text = (shape.text_frame.text or "").strip()
        except Exception:
            continue
        name = str(getattr(shape, "name", "") or "").lower()
        if _is_sample_layout_text(text) or "date" in name:
            try:
                shape.text_frame.clear()
            except Exception:
                try:
                    shape.text_frame.text = ""
                except Exception:
                    pass


def _norm_line(text: str) -> str:
    return " ".join((text or "").lower().split())


def _filter_bullets_against_title(title: str, lines: list[str]) -> list[str]:
    """Drop body lines that repeat the slide title (common planner duplicate)."""
    norm_title = _norm_line(title)
    kept: list[str] = []
    for line in lines:
        raw = str(line or "").strip()
        if not raw:
            continue
        norm = _norm_line(raw)
        if norm_title:
            if norm == norm_title or norm_title in norm or norm in norm_title:
                continue
            if norm.startswith("frame ") and norm_title in _norm_line(raw[6:]):
                continue
        kept.append(raw)
    return kept


def _fill_text_frame(shape, lines: list[str], *, title: str = "") -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    lines = _filter_bullets_against_title(title, lines)
    if not lines:
        tf.text = ""
        return
    tf.text = lines[0][:800]
    for extra in lines[1:8]:
        p = tf.add_paragraph()
        p.text = extra[:800]
        p.level = 0


_TITLE_PH = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}
_BODY_PH = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.VERTICAL_BODY,
}
_SUBTITLE_PH = {
    PP_PLACEHOLDER.SUBTITLE,
}


def _format_deck_date(metadata: dict[str, Any] | None = None) -> str:
    meta = metadata if isinstance(metadata, dict) else {}
    raw = str(meta.get("presentation_date") or meta.get("date") or "").strip()
    if raw:
        return raw[:80]
    return date.today().strftime("%B %d, %Y")


def _fill_placeholder(
    slide,
    *,
    title: str,
    bullets: list[str],
    subtitle: str = "",
    deck_date: str = "",
    regions: dict[str, Any] | None = None,
    skip_body: bool = False,
    prefer_generated_body: bool = False,
) -> None:
    """Populate placeholders XOR generated shapes — never both for the same role.

    Zinnia masters often put an OBJECT content placeholder in the title slot and a
    BODY placeholder below. Prefer real TITLE/BODY types; only then assign leftover
    OBJECT slots (top-most unused → title, next → body). Never fill one shape twice.
    """
    _clear_placeholder_sample_text(slide)
    _clear_layout_sample_text(slide)
    title_ph = None
    subtitle_ph = None
    body_ph = None
    date_ph = None
    object_phs: list[Any] = []
    if slide.shapes.title is not None:
        try:
            if slide.shapes.title.placeholder_format.type in _TITLE_PH:
                title_ph = slide.shapes.title
        except Exception:
            title_ph = slide.shapes.title
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            ph_type = shape.placeholder_format.type
        except ValueError:
            continue
        if ph_type in _TITLE_PH:
            if title_ph is None:
                title_ph = shape
        elif ph_type in _SUBTITLE_PH:
            if subtitle_ph is None:
                subtitle_ph = shape
        elif ph_type in _BODY_PH:
            if body_ph is None:
                body_ph = shape
        elif ph_type == PP_PLACEHOLDER.DATE:
            if date_ph is None:
                date_ph = shape
        elif ph_type == PP_PLACEHOLDER.OBJECT:
            object_phs.append(shape)
    unused = [
        s
        for s in object_phs
        if s is not title_ph and s is not subtitle_ph and s is not body_ph and s is not date_ph
    ]
    unused.sort(key=lambda s: (int(getattr(s, "top", 0) or 0), int(getattr(s, "height", 0) or 0)))
    subtitle_text = (subtitle or "").strip()
    for shape in unused:
        if title_ph is None:
            title_ph = shape
        elif body_ph is None and not subtitle_text:
            body_ph = shape
        elif subtitle_ph is None and subtitle_text:
            subtitle_ph = shape
        elif body_ph is None:
            body_ph = shape
    title_set = False
    subtitle_set = False
    body_set = False
    if title_ph is not None:
        title_ph.text_frame.word_wrap = True
        title_ph.text_frame.text = title[:160]
        title_set = True
    if subtitle_ph is not None and subtitle_text:
        subtitle_ph.text_frame.word_wrap = True
        subtitle_ph.text_frame.text = subtitle_text[:240]
        subtitle_set = True
    if date_ph is not None and deck_date:
        date_ph.text_frame.text = deck_date[:80]
    elif deck_date:
        for shape in slide.shapes:
            if not shape.has_text_frame or _is_placeholder(shape):
                continue
            name = str(getattr(shape, "name", "") or "").lower()
            text = (shape.text_frame.text or "").strip()
            if "date" in name or _is_sample_layout_text(text):
                shape.text_frame.text = deck_date[:80]
                break
    if not skip_body and body_ph is not None and body_ph is not title_ph and not prefer_generated_body:
        filtered = _filter_bullets_against_title(title, bullets)
        if filtered:
            _fill_text_frame(body_ph, filtered, title=title)
            body_set = True
        else:
            try:
                body_ph.text_frame.clear()
            except Exception:
                pass
    if not title_set:
        x, y, cx, cy = _geom_box(regions, "title", (0.5, 0.28, 9.0, 0.7))
        box = slide.shapes.add_textbox(x, y, cx, cy)
        box.text_frame.word_wrap = True
        box.text_frame.text = title[:160]
        title_set = True
    if not subtitle_set and subtitle_text:
        x, y, cx, cy = _geom_box(regions, "subtitle", (0.5, 0.95, 9.0, 0.45))
        box = slide.shapes.add_textbox(x, y, cx, cy)
        box.text_frame.word_wrap = True
        box.text_frame.text = subtitle_text[:240]
        subtitle_set = True
    if skip_body:
        return
    if not body_set and bullets:
        filtered = _filter_bullets_against_title(title, bullets)
        if filtered:
            x, y, cx, cy = _geom_box(regions, "body", (0.5, 1.2, 9.0, 4.8))
            box = slide.shapes.add_textbox(x, y, cx, cy)
            _fill_text_frame(box, filtered, title=title)


def _set_notes(slide, notes: str) -> None:
    text = (notes or "").strip()[:4000]
    if not text:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _prepare_image_assets(plan: dict[str, Any], *, user_id: str) -> None:
    from app.services.mentrix.presentation.asset_resolver import store_example_image

    uid = (user_id or "anon").strip() or "anon"
    for slide_spec in list(plan.get("slides") or []):
        for block in list(slide_spec.get("blocks") or []):
            if str(block.get("kind") or "") != "image":
                continue
            content = block.get("content") if isinstance(block.get("content"), dict) else {}
            if str(content.get("asset_id") or "").strip():
                continue
            meta = store_example_image(user_id=uid, label=str(slide_spec.get("title") or "ZECT")[:32])
            content["asset_id"] = meta["asset_id"]
            block["content"] = content
            block["provenance"] = {
                "source": "example",
                "generated": True,
                "note": "Generated placeholder, not a user photo",
            }
            block["validation"] = {"ok": True, "errors": []}


def validate_generated_pptx(data: bytes, *, n_slides: int) -> dict[str, Any]:
    zf = inspect_pptx_archive(data)
    try:
        names = [i.filename.replace("\\", "/") for i in zf.infolist()]
    finally:
        zf.close()
    if "ppt/presentation.xml" not in names:
        raise UnsafePptxError("missing_presentation_xml")
    slides = parse_pptx_bytes(data)
    if len(slides) < 1:
        raise UnsafePptxError("no_slides")
    if len(slides) != int(n_slides):
        raise UnsafePptxError("slide_count_mismatch")
    return {"ok": True, "slide_count": len(slides), "notes": sum(1 for s in slides if (s.get("notes") or "").strip())}


def render_plan_to_pptx(
    plan: dict[str, Any],
    *,
    template_path: Path | None = None,
    definition: dict[str, Any] | None = None,
    user_id: str = "anon",
) -> bytes:
    """Render a PresentationPlan to PPTX bytes using canonical visual blocks."""
    from app.services.mentrix.presentation.blocks import VISUAL_KINDS, text_lines
    from app.services.mentrix.presentation.layout import choose_layout, place_blocks
    from app.services.mentrix.presentation.visual import paint_block

    slides_in = list(plan.get("slides") or [])
    if not slides_in:
        raise UnsafePptxError("plan_has_no_slides")
    _prepare_image_assets(plan, user_id=user_id)
    prs = None
    if template_path and Path(template_path).is_file():
        raw = Path(template_path).read_bytes()
        zf = inspect_pptx_archive(raw)
        zf.close()
        try:
            prs = Presentation(io.BytesIO(raw))
            _clear_slides(prs)
        except Exception as exc:
            raise UnsafePptxError("template_open_failed") from exc
    if prs is None:
        prs = Presentation()
        cx = (definition or {}).get("slide_size", {}).get("cx") if definition else None
        cy = (definition or {}).get("slide_size", {}).get("cy") if definition else None
        if cx and cy:
            try:
                prs.slide_width = Emu(int(cx))
                prs.slide_height = Emu(int(cy))
            except (TypeError, ValueError):
                pass
    for slide_spec in slides_in:
        layout_intent = choose_layout(slide_spec)
        preferred = str(slide_spec.get("master_layout_name") or "")
        layout = _pick_layout(prs, layout_intent, definition, preferred_name=preferred)
        slide = prs.slides.add_slide(layout)
        try:
            fallback_index = int(slide_spec.get("index") or 0)
        except (TypeError, ValueError):
            fallback_index = 0
        title = str(slide_spec.get("title") or f"Slide {fallback_index + 1}")
        subtitle = str(slide_spec.get("key_message") or "").strip()
        deck_date = _format_deck_date(plan.get("metadata") if isinstance(plan.get("metadata"), dict) else None)
        bullets = text_lines(list(slide_spec.get("blocks") or []))
        if not bullets:
            blocks = list(slide_spec.get("content_blocks") or [])
            bullets = [str(b.get("text") or "").strip() for b in blocks if str(b.get("text") or "").strip()]
        visuals = [b for b in list(slide_spec.get("blocks") or []) if str(b.get("kind") or "") in VISUAL_KINDS]
        skip_body = layout_intent in {"full_image", "quote"}
        regions = slide_spec.get("composed_regions") if isinstance(slide_spec.get("composed_regions"), dict) else None
        if visuals and regions:
            body = regions.get("body") if isinstance(regions.get("body"), dict) else None
            visual = regions.get("visual") if isinstance(regions.get("visual"), dict) else None
            if body and visual and body.get("x") == visual.get("x") and body.get("y") == visual.get("y"):
                skip_body = True
        _fill_placeholder(
            slide,
            title=title,
            subtitle=subtitle,
            deck_date=deck_date if layout_intent in {"title", "title_body", "section", "closing"} else "",
            bullets=[] if skip_body else bullets,
            regions=regions,
            skip_body=skip_body,
            prefer_generated_body=bool(visuals),
        )
        placements = place_blocks(slide_spec, prs=prs, definition=definition)
        for item in placements:
            paint_block(slide, item["block"], item["geometry"], user_id=user_id)
        _set_notes(slide, str(slide_spec.get("notes_intent") or ""))
    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()
    expected = int(plan.get("requested_slide_count") or plan.get("n_slides") or len(slides_in))
    validate_generated_pptx(data, n_slides=expected)
    return data


def write_pptx(data: bytes, dest: Path) -> Path:
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(data)
    return dest


def is_fallback_template_id(template_id: str) -> bool:
    return (template_id or "").strip().lower() in _FALLBACK_IDS
