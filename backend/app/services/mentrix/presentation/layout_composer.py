"""Template-aware layout composition. Uses the selected layout's placeholders, not layout[0]."""

from __future__ import annotations

from typing import Any

from pptx.util import Inches

from app.services.mentrix.presentation.content_capacity import apply_content_budget, dedupe_semantic_blocks
from app.services.mentrix.presentation.content_intent import choose_slide_intent, intent_to_layout
from app.services.mentrix.presentation.quality_policy import SAFE_MARGIN_EMU, TITLE_BODY_GAP_EMU, slide_size_emu


def _layouts(definition: dict[str, Any] | None) -> list[dict[str, Any]]:
    return [item for item in list((definition or {}).get("layouts") or []) if isinstance(item, dict)]


def _ph(layout: dict[str, Any]) -> list[dict[str, Any]]:
    return [p for p in list(layout.get("placeholders") or []) if isinstance(p, dict)]


def _int(value: Any, default: int = 0) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _geom(ph: dict[str, Any]) -> dict[str, int]:
    raw = ph.get("geometry") if isinstance(ph.get("geometry"), dict) else ph
    return {
        "x": _int(raw.get("x")),
        "y": _int(raw.get("y")),
        "cx": max(1, _int(raw.get("cx"), 1)),
        "cy": max(1, _int(raw.get("cy"), 1)),
    }


def _score_layout(
    layout: dict[str, Any],
    *,
    want: str,
    intent: str,
    used: list[str],
    purpose: str = "",
    visual_intent: str = "",
) -> int:
    name = str(layout.get("name") or "").lower()
    purpose_l = (purpose or "").lower()
    visual = (visual_intent or "").lower()
    score = 0
    if purpose_l in {"opening", "title", "intro"} and any(k in name for k in ("title page", "cover", "title only")):
        score += 90
    if purpose_l in {"closing", "summary", "next steps", "ask"} and any(k in name for k in ("title only", "section", "closing")):
        score += 75
    if purpose_l in {"section", "divider"} and "section" in name:
        score += 85
    if visual in {"chart", "table", "image", "diagram"} and any(k in name for k in ("1 column", "subtitle + 1", "blank", "two column")):
        score += 55
    if want == "two_column" and ("2 column" in name or "two column" in name):
        score += 80
    if want == "title_body" and any(k in name for k in ("1 column", "subtitle + 1", "title + 1", "header")):
        score += 70
    if want in {"chart_commentary", "table", "metrics"} and any(k in name for k in ("1 column", "subtitle + 1", "header")):
        score += 60
    if want == "diagram" and ("blank" in name or "1_blank" in name):
        score += 50
    if intent == "TEXT" and ("title page" in name or "cover" in name or "title only" in name):
        score += 75
    if "blank" in name and want != "diagram":
        score -= 20
    if "quote" in name and intent != "TEXT":
        score -= 30
    used_count = used.count(str(layout.get("name") or ""))
    score -= used_count * 25
    if _ph(layout):
        score += 5
    return score


def pick_template_layout(
    definition: dict[str, Any] | None,
    slide: dict[str, Any],
    *,
    used_names: list[str] | None = None,
    exclude_names: list[str] | None = None,
) -> dict[str, Any] | None:
    from app.services.mentrix.presentation.template_semantics import enrich_definition_semantics, region_overlaps_protected

    definition = enrich_definition_semantics(definition)
    layouts = _layouts(definition)
    if not layouts:
        return None
    requested = str(slide.get("master_layout_name") or "").strip()
    excluded = {str(n) for n in list(exclude_names or []) if n}
    if requested and requested not in excluded:
        for layout in layouts:
            if str(layout.get("name") or "") == requested:
                return layout
    intent = str(slide.get("content_intent") or choose_slide_intent(slide))
    want = intent_to_layout(intent)
    used = list(used_names or [])
    purpose = str(slide.get("purpose") or "")
    visual_intent = str(slide.get("visual_intent") or slide.get("visual_choice") or "")

    def layout_score(layout: dict[str, Any]) -> int:
        name = str(layout.get("name") or "")
        if name in excluded:
            return -9999
        score = _score_layout(
            layout,
            want=want,
            intent=intent,
            used=used,
            purpose=purpose,
            visual_intent=visual_intent,
        )
        sem = layout.get("semantic_map") if isinstance(layout.get("semantic_map"), dict) else {}
        protected = list(sem.get("protected_regions") or [])
        body_regions = list(sem.get("body_regions") or [])
        for body in body_regions:
            if region_overlaps_protected(body, protected):
                score -= 120
        tags = set(sem.get("purpose_tags") or [])
        if purpose.lower() in tags:
            score += 40
        return score

    ranked = sorted(layouts, key=layout_score, reverse=True)
    choice = ranked[0] if ranked else layouts[0]
    if layout_score(choice) < -100 and len(ranked) > 1:
        choice = ranked[1]
    rejected = [str(l.get("name") or "") for l in ranked[1:4] if str(l.get("name") or "")]
    slide["layout_selection_reason"] = (
        f"purpose={purpose or 'body'} intent={intent} layout={choice.get('name')} "
        f"score={layout_score(choice)} rejected={','.join(rejected)}"
    )
    slide["layout_rejected_candidates"] = rejected
    return choice


def compose_regions(
    definition: dict[str, Any] | None,
    layout: dict[str, Any] | None,
    *,
    split_visual: bool = False,
) -> dict[str, dict[str, int]]:
    """Title/body/visual regions from the selected layout, respecting protected decoration."""
    from app.services.mentrix.presentation.template_semantics import enrich_definition_semantics

    definition = enrich_definition_semantics(definition)
    cx, cy = slide_size_emu(definition)
    sem = layout.get("semantic_map") if isinstance((layout or {}).get("semantic_map"), dict) else {}
    if sem.get("safe_content_bounds"):
        safe = dict(sem["safe_content_bounds"])
    else:
        m = SAFE_MARGIN_EMU
        safe = {"x": m, "y": m, "cx": cx - 2 * m, "cy": cy - 2 * m}
    placeholders = _ph(layout or {})
    ordered = sorted(placeholders, key=lambda p: (int(p.get("y") or 0), int(p.get("x") or 0)))
    title_ph = next((p for p in ordered if str(p.get("type") or "").upper() in {"TITLE", "CTRTITLE"}), None)
    subtitle_ph = next((p for p in ordered if str(p.get("type") or "").upper() in {"SUBTITLE", "SUBTITLE"}), None)
    body_phs = [p for p in ordered if p is not title_ph and p is not subtitle_ph]
    if title_ph is None and ordered:
        title_ph = ordered[0]
        body_phs = ordered[1:]
    if title_ph is None:
        title = {"x": safe["x"], "y": safe["y"], "cx": safe["cx"], "cy": int(Inches(0.85))}
    else:
        title = _geom(title_ph)
    if body_phs:
        body = _geom(body_phs[0])
        if body["y"] < title["y"] + title["cy"] + TITLE_BODY_GAP_EMU:
            body["y"] = title["y"] + title["cy"] + TITLE_BODY_GAP_EMU
            body["cy"] = max(int(Inches(0.6)), cy - SAFE_MARGIN_EMU - body["y"])
    else:
        body_y = title["y"] + title["cy"] + TITLE_BODY_GAP_EMU
        body = {"x": safe["x"], "y": body_y, "cx": safe["cx"], "cy": max(1, cy - SAFE_MARGIN_EMU - body_y)}
    if subtitle_ph is not None:
        subtitle = _geom(subtitle_ph)
    elif sem.get("subtitle_region"):
        subtitle = dict(sem["subtitle_region"])
    else:
        subtitle_y = title["y"] + title["cy"] + int(TITLE_BODY_GAP_EMU // 2)
        subtitle = {
            "x": safe["x"],
            "y": subtitle_y,
            "cx": safe["cx"],
            "cy": max(int(Inches(0.45)), int(Inches(0.55))),
        }
        body["y"] = max(body["y"], subtitle["y"] + subtitle["cy"] + TITLE_BODY_GAP_EMU)
        body["cy"] = max(int(Inches(0.6)), cy - SAFE_MARGIN_EMU - body["y"])
    visual = dict(body)
    if split_visual and visual.get("x") == body.get("x") and visual.get("y") == body.get("y"):
        gap = int(Inches(0.16))
        text_cy = max(int(Inches(0.9)), body["cy"] // 3)
        visual = {
            "x": body["x"],
            "y": body["y"] + text_cy + gap,
            "cx": body["cx"],
            "cy": max(int(Inches(0.8)), body["cy"] - text_cy - gap),
        }
        body = {**body, "cy": text_cy}
    # Clamp body/visual inside semantic safe bounds when decoration intrudes
    if sem.get("safe_content_bounds"):
        sb = sem["safe_content_bounds"]
        for key, region in (("body", body), ("visual", visual)):
            region["x"] = max(region["x"], sb["x"])
            region["y"] = max(region["y"], sb["y"])
            max_cx = sb["x"] + sb["cx"] - region["x"]
            max_cy = sb["y"] + sb["cy"] - region["y"]
            region["cx"] = min(region["cx"], max(int(Inches(0.8)), max_cx))
            region["cy"] = min(region["cy"], max(int(Inches(0.6)), max_cy))
    # safe_content_bounds above is computed once from the layout's OWN
    # primary body placeholder -- when that placeholder gets reassigned as
    # the title above (a layout with only one placeholder and no
    # dedicated title type) the body/visual regions actually placed here
    # land at a different y than the one safe_content_bounds was shrunk
    # for, so a protected decoration (e.g. a Zinnia layout's own vertical
    # divider bar) that only intrudes at THIS y-band never gets checked at
    # all. Re-shrink the actual final regions directly against every
    # protected region, not just the cached safe box. Reproduced live:
    # the Zinnia "Gradient Bottom" layout's divider still collided with
    # generated body text after only the safe_content_bounds fix.
    protected = list(sem.get("protected_regions") or [])
    if protected:
        from app.services.mentrix.presentation.template_semantics import shrink_region_away_from

        for key, region in (("body", body), ("visual", visual)):
            for prot in protected:
                region.update(shrink_region_away_from(region, prot, slide_cx=cx, slide_cy=cy))
    return {"title": title, "subtitle": subtitle, "body": body, "visual": visual, "safe": safe, "semantic_map": sem}


def compose_plan(plan: dict[str, Any], definition: dict[str, Any] | None, *, prompt: str = "") -> dict[str, Any]:
    from app.services.mentrix.presentation.template_semantics import enrich_definition_semantics

    definition = enrich_definition_semantics(definition)
    used: list[str] = []
    from app.services.mentrix.presentation.blocks import VISUAL_KINDS

    for slide in list(plan.get("slides") or []):
        intent = choose_slide_intent(slide, purpose=str(slide.get("purpose") or ""), prompt=prompt)
        slide["content_intent"] = intent
        exclude = list(slide.get("_layout_exclude") or [])
        layout = pick_template_layout(definition, slide, used_names=used, exclude_names=exclude)
        name = str((layout or {}).get("name") or "")
        if name:
            slide["master_layout_name"] = name
            used.append(name)
        has_visual = any(str(b.get("kind") or "") in VISUAL_KINDS for b in list(slide.get("blocks") or []))
        regions = compose_regions(definition, layout, split_visual=has_visual)
        slide["composed_regions"] = regions
        apply_content_budget(slide, regions)
        dedupe_semantic_blocks(slide)
        slide["layout"] = intent_to_layout(intent)
    return plan
