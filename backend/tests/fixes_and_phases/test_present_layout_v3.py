"""V3 layout generation closure — count semantics, template semantics, quality unity."""

from __future__ import annotations

from app.services.mentrix.presentation import prepare_prompt_deck
from app.services.mentrix.presentation.audience import prompt_adapter
from app.services.mentrix.presentation.layout_composer import compose_plan, pick_template_layout
from app.services.mentrix.presentation.quality_critic import critique_plan
from app.services.mentrix.presentation.quality_repair import apply_repairs
from app.services.mentrix.presentation.rendered_quality import inspect_rendered_document
from app.services.mentrix.presentation.template_semantics import (
    build_layout_semantic_map,
    classify_decorative_shape,
    enrich_definition_semantics,
    region_overlaps_protected,
)


def test_explicit_count_removes_audience_hint():
    text = prompt_adapter("general", "AI Agentic vs Graph", requested_slide_count=3)
    assert "Target: 3 slides" in text
    assert "Target ~6" not in text


def test_outline_exact_requested_count():
    out = prepare_prompt_deck(
        prompt="Difference between AI Agentic and the Graph, loop and KV cache",
        audience_id="general",
        requested_slide_count=3,
    )
    assert out["requested_slide_count"] == 3
    assert out["n_slides_hint"] == 3
    assert "Target: 3 slides" in out["adapted_prompt"]
    assert "Target ~6" not in out["adapted_prompt"]


def test_decoration_not_content_region():
    role = classify_decorative_shape(
        geometry={"x": 4000000, "y": 500000, "cx": 800000, "cy": 5000000},
        name="Rectangle 11",
        slide_cx=12192000,
        slide_cy=6858000,
    )
    assert role in {"PROTECTED_BRAND_ELEMENT", "LAYOUT_DECORATION"}


def test_template_shape_semantic_classification():
    layout = {
        "name": "Title and Content",
        "placeholders": [{"type": "title", "geometry": {"x": 100, "y": 100, "cx": 800, "cy": 400}}],
        "shapes": [{"role": "PROTECTED_BRAND_ELEMENT", "geometry": {"x": 9000000, "y": 0, "cx": 2000000, "cy": 6858000}}],
    }
    sem = build_layout_semantic_map(layout, slide_cx=12192000, slide_cy=6858000)
    assert sem["protected_regions"]
    assert sem["safe_content_bounds"]["cx"] < 12192000


def test_layout_planner_purpose_driven():
    definition = enrich_definition_semantics(
        {
            "slide_size": {"cx": 12192000, "cy": 6858000},
            "layouts": [
                {"name": "Title Slide", "placeholders": [{"type": "title", "geometry": {"x": 1, "y": 1, "cx": 100, "cy": 100}}]},
                {"name": "Title and Content", "placeholders": [{"type": "body", "geometry": {"x": 1, "y": 200, "cx": 500, "cy": 400}}]},
            ],
        }
    )
    slide = {"purpose": "opening", "title": "Intro", "content_blocks": [{"text": "Hello"}]}
    layout = pick_template_layout(definition, slide)
    assert "Title" in str(layout.get("name") or "")


def test_protected_template_region():
    body = {"x": 100, "y": 100, "cx": 500, "cy": 400}
    prot = {"x": 200, "y": 150, "cx": 300, "cy": 200}
    assert region_overlaps_protected(body, [prot])


def test_rendered_overlap_detected():
    doc = {
        "slide_cx": 12192000,
        "slide_cy": 6858000,
        "slides": [
            {
                "blocks": [
                    {"kind": "text", "content": {"text": "A" * 80}, "geometry": {"x": 100, "y": 100, "cx": 400, "cy": 200}},
                    {"kind": "text", "content": {"text": "B" * 80}, "geometry": {"x": 120, "y": 120, "cx": 400, "cy": 200}},
                ]
            }
        ],
    }
    out = inspect_rendered_document(doc)
    assert out["rendered_overlap_count"] >= 1
    assert out["status"] == "FAIL"


def test_quality_summary_matches_rendered_findings():
    from app.services.mentrix.presentation.deck_catalog import quality_gate_for_path
    from app.services.mentrix.presentation.native_provider import _unique_pptx_path
    from app.services.pptx_paths import default_pptx_save_dir
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Overlap test"
    dest = _unique_pptx_path(default_pptx_save_dir(), "v3-quality-gate.pptx")
    prs.save(str(dest))
    gate = quality_gate_for_path(str(dest))
    if gate.get("rendered_overlap_count", 0) > 0:
        assert gate["overlap_count"] >= gate["rendered_overlap_count"]


def test_repair_switches_layout():
    plan = {
        "requested_slide_count": 3,
        "n_slides": 3,
        "slides": [
            {
                "title": "One",
                "master_layout_name": "Layout A",
                "content_blocks": [{"text": "a"}],
                "composed_regions": {"body": {"x": 1, "y": 1, "cx": 1, "cy": 1}},
            }
        ],
    }
    report = {
        "slides": [{"repairs": ["change_layout"], "findings": ["overlap"]}],
    }
    apply_repairs(plan, report)
    assert plan["slides"][0].get("master_layout_name") is None
    assert "Layout A" in plan["slides"][0].get("_layout_exclude", [])


def test_layout_capacity_rejects_overflow():
    slide = {
        "title": "Long",
        "content_blocks": [{"text": "x" * 500} for _ in range(12)],
        "composed_regions": {"body": {"x": 100, "y": 200, "cx": 200, "cy": 100}},
    }
    definition = {"slide_size": {"cx": 12192000, "cy": 6858000}, "layouts": []}
    row = critique_plan({"slides": [slide]}, definition, prompt="test")["slides"][0]
    assert row["status"] in {"FAIL", "REPAIRABLE"}
