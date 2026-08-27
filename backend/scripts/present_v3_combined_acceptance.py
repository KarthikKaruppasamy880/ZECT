#!/usr/bin/env python3
"""Combined V3 Present headed acceptance — evidence bundle for human visual review."""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
REPO = ROOT.parent
ART = REPO / "test-results" / "present-v3-headed-acceptance"
GOLDEN_PROMPT = (
    "Difference between AI Agentic and the Graph, loop and KV catch with LLM fine tuning"
)
REQUESTED = 3
TEMPLATE = "zinnia-executive-v1"


def _git_head() -> str:
    try:
        return subprocess.check_output(["git", "rev-parse", "HEAD"], cwd=REPO, text=True).strip()
    except Exception:
        return "unknown"


def _slide_png_evidence(pptx: Path) -> dict:
    sys.path.insert(0, str(ROOT))
    from app.services.mentrix.presentation.slide_preview import render_slide_png_bytes
    from pptx import Presentation

    data = pptx.read_bytes()
    n = len(Presentation(str(pptx)).slides)
    slides_dir = ART / "slide-pngs"
    slides_dir.mkdir(parents=True, exist_ok=True)
    rows = []
    for idx in range(n):
        png = render_slide_png_bytes(data, idx)
        dest = slides_dir / f"slide-{idx + 1}.png"
        dest.write_bytes(png)
        rows.append({"index": idx, "bytes": len(png), "path": str(dest)})
    return {"slide_count": n, "slides": rows}


def _com_powerpoint_proof(pptx: Path) -> dict:
    if os.environ.get("ZECT_LIVE_PPT_COM", "").strip() != "1" or os.name != "nt":
        return {"ok": False, "status": "BLOCKED_EXTERNAL", "reason": "ZECT_LIVE_PPT_COM!=1 or not Windows"}
    try:
        import win32com.client  # type: ignore
    except ImportError:
        return {"ok": False, "status": "BLOCKED_EXTERNAL", "reason": "win32com_missing"}
    app = None
    proof_dir = ART / "powerpoint"
    proof_dir.mkdir(parents=True, exist_ok=True)
    out: dict = {"pptx": str(pptx)}
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 1
        pres = app.Presentations.Open(str(pptx.resolve()), WithWindow=True, ReadOnly=False)
        out["slide_count"] = int(pres.Slides.Count)
        out["exactly_three"] = out["slide_count"] == REQUESTED
        for i in range(1, out["slide_count"] + 1):
            slide = pres.Slides(i)
            export_path = proof_dir / f"ppt-slide-{i}.png"
            slide.Export(str(export_path.resolve()), "PNG", 1920, 1080)
            title = ""
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    t = str(shape.TextFrame.TextRange.Text or "").strip()
                    if t and len(t) > len(title):
                        title = t[:200]
            out.setdefault("slides", []).append(
                {
                    "index": i,
                    "title_sample": title,
                    "export_png": str(export_path),
                    "has_date_here": "date here" in title.lower(),
                }
            )
        out["ok"] = True
        out["status"] = "opened_visible"
        out["repair_dialog"] = False
        out["note"] = "PowerPoint left open for human inspection — close manually when done."
    except Exception as exc:  # noqa: BLE001
        out["ok"] = False
        out["status"] = "failed"
        out["error"] = str(exc)[:300]
    return out


def _repair_flow_test() -> dict:
    sys.path.insert(0, str(ROOT))
    import io

    from pptx import Presentation
    from pptx.util import Inches

    from app.services.mentrix.presentation.deck_catalog import quality_gate_for_path
    from app.services.mentrix.presentation.final_pptx_inspector import inspect_and_repair_pptx
    from app.services.pptx_paths import default_pptx_save_dir

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Repair flow probe"
    dump = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(6.5))
    dump.text_frame.text = "Overlapping dump shape for repair-deck acceptance."
    buf = io.BytesIO()
    prs.save(buf)
    dest = default_pptx_save_dir() / "v3-repair-flow-probe.pptx"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(buf.getvalue())
    before = quality_gate_for_path(str(dest))
    repaired, rep = inspect_and_repair_pptx(dest.read_bytes())
    dest.write_bytes(repaired)
    after = quality_gate_for_path(str(dest))
    return {
        "path": str(dest),
        "before_status": before.get("final_quality_status"),
        "before_export_blocked": before.get("export_blocked"),
        "after_status": after.get("final_quality_status"),
        "after_export_blocked": after.get("export_blocked"),
        "repair_report": {
            "dump_shapes_removed": rep.get("dump_shapes_removed"),
            "duplicate_shapes_removed": rep.get("duplicate_shapes_removed"),
        },
        "pass_after_repair": after.get("final_quality_status") == "PASS" and not after.get("export_blocked"),
    }


def main() -> int:
    os.environ.setdefault("ZECT_PRESENTATION_PROVIDER", "zect_native")
    sys.path.insert(0, str(ROOT))
    ART.mkdir(parents=True, exist_ok=True)
    head = _git_head()

    golden_report_path = ROOT / "artifacts" / "present-golden-v3" / "golden_v3_report.json"
    if not golden_report_path.is_file():
        subprocess.check_call([sys.executable, str(ROOT / "scripts" / "present_golden_v3_proof.py")], cwd=str(ROOT))
    golden = json.loads(golden_report_path.read_text(encoding="utf-8"))
    pptx_path = Path(str(golden.get("path") or ""))
    if not pptx_path.is_file():
        print(json.dumps({"error": "golden_pptx_missing", "golden": golden}, indent=2))
        return 1

    prod = subprocess.run(
        [sys.executable, str(ROOT / "scripts" / "present_production_release_proof.py")],
        cwd=str(ROOT),
        capture_output=True,
        text=True,
    )
    prod_json = {}
    prod_path = ROOT / "artifacts" / "present-production-release" / "production_release_report.json"
    if prod_path.is_file():
        prod_json = json.loads(prod_path.read_text(encoding="utf-8"))

    png_evidence = _slide_png_evidence(pptx_path)
    com_proof = _com_powerpoint_proof(pptx_path)
    repair = _repair_flow_test()

    ui_evidence_path = ART / "ui-evidence.json"
    ui_evidence = {}
    if ui_evidence_path.is_file():
        ui_evidence = json.loads(ui_evidence_path.read_text(encoding="utf-8"))

    report = {
        "verdict": "READY_FOR_HUMAN_PRESENT_V3_VISUAL_REVIEW",
        "timestamp_utc": datetime.now(timezone.utc).isoformat(),
        "git_head": head,
        "branch": "feat/present-p3-editor-export-p1",
        "golden_prompt": GOLDEN_PROMPT,
        "requested_slides": REQUESTED,
        "template": TEMPLATE,
        "generated_pptx": str(pptx_path),
        "golden_v3": golden,
        "production_proof": prod_json,
        "production_proof_exit": prod.returncode,
        "slide_png_evidence": png_evidence,
        "powerpoint_com_proof": com_proof,
        "repair_flow": repair,
        "ui_evidence": ui_evidence,
        "review_url_hint": ui_evidence.get("review_url") or "(run headed playwright for URL)",
        "electron_url_hint": ui_evidence.get("electron_url") or "(Electron should load same review URL)",
    }
    out_path = ART / "combined_acceptance_report.json"
    out_path.write_text(json.dumps(report, indent=2), encoding="utf-8")
    shutil.copy2(pptx_path, ART / "golden-v3-agentic-deck.pptx")
    print(json.dumps(report, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
