"""Editable PPTX charts via python-pptx (Presenton-parity types + SVG fallback)."""

from __future__ import annotations

from typing import Any

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu, Inches

from app.services.mentrix.presentation.blocks import CHART_TYPES

# Presenton labels map to these ids. polar/progress/gauge have no OOXML type.
_XL_NATIVE = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "radar": XL_CHART_TYPE.RADAR,
    "area": XL_CHART_TYPE.AREA,
    "stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_horizontal": getattr(XL_CHART_TYPE, "BAR_STACKED", XL_CHART_TYPE.BAR_CLUSTERED),
    "scatter": getattr(XL_CHART_TYPE, "XY_SCATTER", XL_CHART_TYPE.LINE),
}

_OOXML_ALIAS = {
    "polar": "radar",
    "progress": "column",
    "gauge": "donut",
}

_SVG_PREFERRED = frozenset({"progress", "gauge"})


def resolve_ooxml_chart_type(chart_type: str) -> str:
    raw = (chart_type or "column").lower().strip()
    if raw not in CHART_TYPES:
        return "column"
    return _OOXML_ALIAS.get(raw, raw)


def chart_data_from_block(block: dict[str, Any]) -> CategoryChartData | None:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    categories = [str(c) for c in list(content.get("categories") or []) if str(c).strip()]
    series = [s for s in list(content.get("series") or []) if isinstance(s, dict)]
    if len(categories) < 2 or not series:
        return None
    data = CategoryChartData()
    data.categories = categories
    width = len(categories)
    for spec in series[:6]:
        values = list(spec.get("values") or [])[:width]
        if len(values) != width:
            return None
        data.add_series(str(spec.get("name") or "Series")[:80], tuple(float(v) for v in values))
    return data


def xy_data_from_block(block: dict[str, Any]) -> XyChartData | None:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    categories = [str(c) for c in list(content.get("categories") or [])]
    series = [s for s in list(content.get("series") or []) if isinstance(s, dict)]
    if not series:
        return None
    data = XyChartData()
    for spec in series[:6]:
        values = list(spec.get("values") or [])
        if len(values) < 2:
            return None
        series_obj = data.add_series(str(spec.get("name") or "Series")[:80])
        for i, val in enumerate(values):
            try:
                x = float(categories[i]) if i < len(categories) else float(i)
            except (TypeError, ValueError):
                x = float(i)
            series_obj.add_data_point(x, float(val))
    return data


def _drop_charts(slide) -> None:
    doomed = []
    for shape in slide.shapes:
        try:
            if getattr(shape, "has_chart", False):
                doomed.append(shape)
        except (ValueError, AttributeError):
            continue
    for shape in doomed:
        el = shape._element  # noqa: SLF001 — python-pptx has no public chart delete
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


def add_chart(slide, block: dict[str, Any], geometry: dict[str, int]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    requested = str(content.get("chart_type") or "column").lower()
    if requested in _SVG_PREFERRED:
        return False
    chart_type = resolve_ooxml_chart_type(requested)
    if chart_type not in _XL_NATIVE:
        chart_type = "column"
    x, y, cx, cy = (Emu(int(geometry[k])) for k in ("x", "y", "cx", "cy"))
    if chart_type == "scatter":
        data = xy_data_from_block(block)
        if data is None:
            return False
        chart = slide.shapes.add_chart(_XL_NATIVE[chart_type], x, y, cx, cy, data).chart
    else:
        data = chart_data_from_block(block)
        if data is None:
            return False
        chart = slide.shapes.add_chart(_XL_NATIVE[chart_type], x, y, cx, cy, data).chart
    title = str(content.get("title") or "").strip()[:160]
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    chart.has_legend = bool(content.get("legend", True))
    return True


def replace_chart_data(slide, block: dict[str, Any]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    requested = str(content.get("chart_type") or "column").lower()
    if requested in _SVG_PREFERRED:
        return False
    wanted = resolve_ooxml_chart_type(requested)
    data = xy_data_from_block(block) if wanted == "scatter" else chart_data_from_block(block)
    if data is None:
        return False
    existing = None
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            existing = shape
            break
    if existing is None:
        return False
    current_xl = None
    try:
        current_xl = existing.chart.chart_type
    except Exception:
        current_xl = None
    wanted_xl = _XL_NATIVE.get(wanted)
    if wanted == "scatter" or (wanted_xl is not None and current_xl is not None and current_xl != wanted_xl):
        geom = {
            "x": int(getattr(existing, "left", Inches(0.6))),
            "y": int(getattr(existing, "top", Inches(1.6))),
            "cx": int(getattr(existing, "width", Inches(8.8))),
            "cy": int(getattr(existing, "height", Inches(4.6))),
        }
        _drop_charts(slide)
        return add_chart(slide, block, geom)
    try:
        existing.chart.replace_data(data)
    except Exception:
        geom = {
            "x": int(getattr(existing, "left", Inches(0.6))),
            "y": int(getattr(existing, "top", Inches(1.6))),
            "cx": int(getattr(existing, "width", Inches(8.8))),
            "cy": int(getattr(existing, "height", Inches(4.6))),
        }
        _drop_charts(slide)
        return add_chart(slide, block, geom)
    title = str(content.get("title") or "").strip()[:160]
    if title:
        try:
            existing.chart.has_title = True
            existing.chart.chart_title.text_frame.text = title
        except Exception:
            pass
    return True
