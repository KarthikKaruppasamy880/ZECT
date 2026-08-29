"""PresentationDocument canvas parse, AI tree patches, save, and export gates."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from app.services.mentrix.presentation.document_io import (
    apply_document_to_pptx,
    powerpoint_open_without_repair,
    validate_export_document,
)
from app.services.mentrix.presentation.presenter_intelligence import grounded_slide_script
from app.services.mentrix.presentation.slide_ai import patch_slide_from_prompt
from app.services.pptx_parse import parse_pptx_bytes
from tests.fixes_and_phases.pptx_fixtures import make_chart_pptx_bytes, make_group_pptx_bytes


def test_parse_chart_includes_series_not_empty_hitbox():
    slides = parse_pptx_bytes(make_chart_pptx_bytes())
    charts = [b for b in (slides[0].get("blocks") or []) if b.get("kind") == "chart"]
    assert charts
    content = charts[0].get("content") or {}
    assert content.get("categories")
    assert content.get("series")
    assert (content["series"][0].get("values") or [])[0] == 1.0


def test_group_child_stays_offset_and_placeholder_dropped():
    slides = parse_pptx_bytes(make_group_pptx_bytes())
    blocks = slides[0].get("blocks") or []
    texts = [b for b in blocks if "HelloGroup" in str((b.get("content") or {}).get("text") or "")]
    assert texts
    geo = texts[0]["geometry"]
    assert geo["x"] == 1_200_000
    assert geo["y"] == 600_000
    assert not any(int((b.get("geometry") or {}).get("cy") or 0) == 800_000 for b in blocks)


def test_slide_ai_document_tree_from_bullets_does_not_invent_kpis():
    blocks = [{"kind": "text", "content": {"text": "Delivery\nRisk\nAsk"}}]
    out = patch_slide_from_prompt(prompt="Turn these bullets into a diagram.", blocks=blocks)
    assert out["ok"] is True
    assert out["action"] == "bullets_to_diagram"
    kinds = [b.get("kind") for b in out["blocks"]]
    assert "diagram" in kinds
    diagram = next(b for b in out["blocks"] if b.get("kind") == "diagram")
    assert "Delivery" in (diagram.get("content") or {}).get("nodes")
    assert "42%" not in str(out)


def test_slide_ai_table_and_density():
    blocks = [{"kind": "text", "content": {"text": "One\nTwo\nThree\nFour"}}]
    table = patch_slide_from_prompt(prompt="Add a comparison table", blocks=blocks, slide_text="One\nTwo")
    assert table["action"] == "add_table"
    dense = patch_slide_from_prompt(prompt="Reduce density — keep the first three points.", blocks=blocks)
    assert dense["action"] == "reduce_density"
    assert dense["text"].count("\n") <= 2


def test_slide_ai_refuses_diagram_without_points():
    out = patch_slide_from_prompt(prompt="Turn these bullets into a diagram.", blocks=[])
    assert out["ok"] is False
    assert out["error"] == "unparsed_prompt"


def test_named_text_round_trip(tmp_path):
    path = tmp_path / "named.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Original"
        title_name = slide.shapes.title.name
    else:
        title_name = ""
    prs.save(str(path))
    apply_document_to_pptx(
        path,
        [
            {
                "index": 0,
                "text": "Should not dump",
                "notes": "Speaker",
                "blocks": [
                    {
                        "kind": "text",
                        "content": {"text": "Updated title", "shape_name": title_name},
                    }
                ],
            }
        ],
    )
    opened = Presentation(str(path))
    assert "Updated title" in (opened.slides[0].shapes.title.text if opened.slides[0].shapes.title else "")
    parsed = parse_pptx_bytes(path.read_bytes())
    assert "speaker" in (parsed[0].get("notes") or "").lower()


def test_validate_export_document_ok(tmp_path):
    path = tmp_path / "export.pptx"
    path.write_bytes(make_chart_pptx_bytes())
    out = validate_export_document(path, expected_slides=1)
    assert out["ok"] is True
    assert out["zip_ok"] is True
    assert out["slide_count"] == 1
    assert out["has_chart"] is True


def test_powerpoint_com_is_blocked_without_flag(monkeypatch, tmp_path):
    monkeypatch.delenv("ZECT_LIVE_PPT_COM", raising=False)
    path = tmp_path / "x.pptx"
    path.write_bytes(make_chart_pptx_bytes())
    out = powerpoint_open_without_repair(path)
    assert out["ok"] is False
    assert out["status"] == "BLOCKED_EXTERNAL"


def test_presenter_uses_block_visuals_not_invented_numbers():
    script = grounded_slide_script(
        {
            "notes": "",
            "text": "",
            "blocks": [
                {"kind": "text", "content": {"text": "Coverage is on track"}},
                {"kind": "chart", "content": {"title": "Coverage"}},
            ],
        },
        slide_index=0,
        slide_count=1,
    )
    assert "Coverage is on track" in script
    assert "chart" in script.lower()
    assert "99%" not in script


def test_image_data_url_round_trip_does_not_drop_pictures(tmp_path):
    from PIL import Image

    png = tmp_path / "dot.png"
    Image.new("RGB", (8, 8), (16, 80, 72)).save(png)
    path = tmp_path / "pic.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(png), Inches(1), Inches(1), Inches(2), Inches(2))
    prs.save(str(path))
    slides = parse_pptx_bytes(path.read_bytes())
    images = [b for b in (slides[0].get("blocks") or []) if b.get("kind") == "image"]
    assert images
    assert str((images[0].get("content") or {}).get("data_url") or "").startswith("data:image/")
    apply_document_to_pptx(path, slides)
    again = parse_pptx_bytes(path.read_bytes())
    assert any(b.get("kind") == "image" for b in (again[0].get("blocks") or []))
