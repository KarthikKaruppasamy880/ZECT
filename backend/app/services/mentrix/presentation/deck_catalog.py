"""Recent generated decks in the allowlisted PPTX output directory."""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.services.mentrix.presentation.native_provider import _unique_pptx_path
from app.services.pptx_paths import default_pptx_save_dir, notes_sidecar_for_pptx, resolve_allowlisted_pptx


def delete_deck(path_str: str) -> dict[str, Any]:
    pptx = resolve_allowlisted_pptx(path_str)
    sidecar = notes_sidecar_for_pptx(pptx)
    pptx.unlink()
    if sidecar.is_file():
        sidecar.unlink()
    return {"ok": True, "path": str(pptx), "deleted": True}


def duplicate_deck(path_str: str) -> Path:
    pptx = resolve_allowlisted_pptx(path_str)
    dest = _unique_pptx_path(default_pptx_save_dir(), pptx.name)
    dest.write_bytes(pptx.read_bytes())
    sidecar = notes_sidecar_for_pptx(pptx)
    if sidecar.is_file():
        notes_sidecar_for_pptx(dest).write_bytes(sidecar.read_bytes())
    return dest


def list_recent_decks(*, limit: int = 24) -> list[dict[str, Any]]:
    root = default_pptx_save_dir()
    files = sorted(root.glob("*.pptx"), key=lambda p: p.stat().st_mtime, reverse=True)
    out: list[dict[str, Any]] = []
    for path in files[: max(1, min(80, limit))]:
        sidecar = notes_sidecar_for_pptx(path)
        n_slides = 0
        try:
            prs = Presentation(str(path))
            n_slides = len(prs.slides)
        except Exception:
            n_slides = 0
        out.append(
            {
                "id": path.stem,
                "name": path.name,
                "path": str(path),
                "modified": int(path.stat().st_mtime),
                "bytes": path.stat().st_size,
                "slide_count": n_slides,
                "has_notes": sidecar.is_file(),
                "preview_available": True,
            }
        )
    return out


def create_blank_pptx(*, filename: str = "untitled.pptx", layout: str = "title_slide") -> Path:
    from app.services.mentrix.presentation.blank_document import write_blank_sidecar

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Untitled presentation"
    dest = _unique_pptx_path(default_pptx_save_dir(), filename)
    buf = io.BytesIO()
    prs.save(buf)
    dest.write_bytes(buf.getvalue())
    write_blank_sidecar(dest, layout=layout)  # type: ignore[arg-type]
    return dest


def instantiate_from_template(template_id: str, user_id: str | int | None = None) -> Path:
    """Copy a template master into the allowlisted deck dir so the editor can open it."""
    from app.services.mentrix.presentation import template_registry as tmpl

    zid = tmpl.canonical_id(template_id) or (template_id or "").strip()
    src = tmpl.source_pptx_path(zid, user_id)
    if src is None and zid in tmpl.zinnia_canonical_ids():
        src = tmpl.repo_zinnia_master_pptx()
    if src is None or not src.is_file():
        raise FileNotFoundError("template_master_missing")
    dest = _unique_pptx_path(default_pptx_save_dir(), f"{zid}.pptx")
    dest.write_bytes(src.read_bytes())
    return dest


def import_pptx_bytes(data: bytes, *, filename: str) -> Path:
    """Copy an uploaded PPTX into the allowlisted deck dir after archive + OOXML checks."""
    from app.services.mentrix.presentation.template_importer import UnsafePptxError, inspect_pptx_archive

    with inspect_pptx_archive(data) as zf:
        names = {info.filename.replace("\\", "/") for info in zf.infolist()}
        if "[Content_Types].xml" not in names or "ppt/presentation.xml" not in names:
            raise UnsafePptxError("not_a_pptx")
    try:
        Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — invalid OOXML must fail closed
        raise UnsafePptxError("unreadable_pptx") from exc
    dest = _unique_pptx_path(default_pptx_save_dir(), filename or "imported.pptx")
    dest.write_bytes(data)
    return dest


def quality_gate_for_path(path_str: str) -> dict[str, Any]:
    from app.services.mentrix.presentation.final_pptx_inspector import inspect_pptx_bytes

    pptx = resolve_allowlisted_pptx(path_str)
    try:
        report = inspect_pptx_bytes(pptx.read_bytes())
    except Exception as exc:  # noqa: BLE001 — fail closed; never 500 a hang-looking export UI
        return {
            "ok": False,
            "path": str(pptx),
            "export_blocked": True,
            "hard_blocked": True,
            "accept_warnings_allowed": False,
            "warnings": [],
            "hard_findings": ["inspect_failed"],
            "quality_passed": False,
            "slide_count": 0,
            "overlap_count": 0,
            "clipped_text_count": 0,
            "covering_dump_count": 0,
            "broken_rel_count": 0,
            "final_quality_status": "FAIL",
            "inspector": {"ok": False, "error": str(exc)[:200], "status": "FAIL"},
        }
    hard = list(report.get("hard_findings") or [])
    blocked = bool(report.get("export_blocked") or hard or report.get("status") == "FAIL")
    warnings: list[str] = []
    if not report.get("has_notes"):
        warnings.append("notes_missing")
    critic: dict[str, Any] = {}
    critic_blocked = False
    try:
        from app.services.mentrix.presentation.document import document_from_pptx_bytes
        from app.services.mentrix.presentation.quality_critic import critique_document

        critic = critique_document(document_from_pptx_bytes(pptx.read_bytes(), path=str(pptx)))
        critic_status = str(critic.get("deck_status") or critic.get("status") or "").upper()
        if critic_status == "FAIL" or critic.get("export_blocked"):
            critic_blocked = True
    except Exception:
        critic = {"ok": False, "error": "document_critic_unavailable"}
    if critic_blocked:
        blocked = True
        if "quality_critic_fail" not in hard:
            hard.append("quality_critic_fail")
        if int(critic.get("duplicate_block_id_count") or 0) > 0 and "duplicate_block_id" not in hard:
            hard.append("duplicate_block_id")
        rendered = critic.get("rendered_quality") if isinstance(critic.get("rendered_quality"), dict) else {}
        if int(rendered.get("rendered_overlap_count") or 0) > 0 and "rendered_overlap" not in hard:
            hard.append("rendered_overlap")
        if int(rendered.get("rendered_clipped_text_count") or 0) > 0 and "rendered_clipping" not in hard:
            hard.append("rendered_clipping")
    from app.services.mentrix.presentation.quality_verdict import merge_quality_reports

    rendered = critic.get("rendered_quality") if isinstance(critic.get("rendered_quality"), dict) else {}
    doc_overlap = int(critic.get("document_overlap_count") or 0)
    rendered_overlap = int(rendered.get("rendered_overlap_count") or 0)
    rendered_clipped = int(rendered.get("rendered_clipped_text_count") or 0)
    inspector_overlap = int(report.get("overlap_count") or 0)
    inspector_clipped = int(report.get("out_of_bounds_count") or 0)
    overlap_total = max(inspector_overlap, doc_overlap, rendered_overlap)
    clipped_total = max(inspector_clipped, rendered_clipped)
    template_conflicts = int(rendered.get("template_conflict_count") or 0)
    near_empty = int(critic.get("near_empty_slide_count") or 0)
    merged = merge_quality_reports(inspector=report, document_critic=critic)
    verdict = merged["final_quality_status"]
    blocked = bool(blocked or merged.get("export_blocked"))
    for item in list(merged.get("hard_findings") or []):
        if item not in hard:
            hard.append(item)
    return {
        "ok": True,
        "path": str(pptx),
        "export_blocked": blocked,
        "hard_blocked": blocked,
        "accept_warnings_allowed": bool(warnings) and not blocked,
        "warnings": warnings,
        "hard_findings": hard,
        "quality_passed": verdict == "PASS",
        "slide_count": report.get("slide_count") or 0,
        "overlap_count": overlap_total,
        "clipped_text_count": clipped_total,
        "rendered_overlap_count": rendered_overlap,
        "rendered_clipped_count": rendered_clipped,
        "template_conflict_count": template_conflicts,
        "near_empty_slide_count": near_empty,
        "duplicate_block_id_count": int(critic.get("duplicate_block_id_count") or 0),
        "covering_dump_count": report.get("covering_dump_count") or 0,
        "broken_rel_count": report.get("broken_rel_count") or 0,
        "final_quality_status": verdict,
        "quality_subchecks": merged.get("subchecks") or {},
        "inspector": report,
        "document_critic": critic,
    }
