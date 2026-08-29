"""Write editor text/notes back into OOXML (S5 round-trip). Sidecar remains source of truth if this fails."""

from __future__ import annotations

import json
import os
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Emu, Inches

from app.services.mentrix.presentation.renderer import _set_notes, validate_generated_pptx
from app.services.pptx_parse import parse_pptx_bytes


def apply_document_to_pptx(
    path: str | Path,
    slides: list[dict[str, Any]],
    *,
    user_id: str = "anon",
) -> dict[str, Any]:
    pptx = Path(path)
    prs = Presentation(str(pptx))
    for spec in slides or []:
        try:
            idx = int(spec.get("index", -1))
        except (TypeError, ValueError):
            continue
        if idx < 0 or idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        named = 0
        for block in spec.get("blocks") or []:
            if isinstance(block, dict) and _apply_named_block(slide, block):
                named += 1
        if named == 0:
            text = _text_from_spec(spec)
            if text:
                _write_slide_text(slide, text)
        _set_notes(slide, str(spec.get("notes") or ""))
        _apply_visual_blocks(slide, spec, user_id=user_id)
    tmp = pptx.with_name(f"{pptx.stem}.zect-tmp{pptx.suffix}")
    try:
        prs.save(str(tmp))
        data = tmp.read_bytes()
        validate_generated_pptx(data, n_slides=len(prs.slides))
        os.replace(tmp, pptx)
        try:
            from app.services.mentrix.presentation.slide_preview import invalidate_slide_previews

            invalidate_slide_previews(pptx)
        except Exception:
            pass
        try:
            from app.services.mentrix.presentation.final_pptx_inspector import inspect_and_repair_pptx

            repaired, _rep = inspect_and_repair_pptx(pptx.read_bytes())
            pptx.write_bytes(repaired)
        except Exception:
            pass
    finally:
        if tmp.exists():
            tmp.unlink(missing_ok=True)
    return {"ok": True, "path": str(pptx), "slide_count": len(prs.slides), "ooxml_roundtrip": True}


def validate_export_document(path: str | Path, *, expected_slides: int | None = None) -> dict[str, Any]:
    """OOXML export gate: zip, slide count, media/chart/table presence. Not a COM visual PASS."""
    pptx = Path(path)
    data = pptx.read_bytes()
    if data[:2] != b"PK":
        return {"ok": False, "error": "not_zip"}
    from app.services.mentrix.presentation.document import inspect_pptx_visuals

    parsed = parse_pptx_bytes(data)
    n = expected_slides if expected_slides is not None else len(parsed)
    meta = validate_generated_pptx(data, n_slides=n)
    visuals = inspect_pptx_visuals(data)
    return {"ok": True, **meta, **visuals, "zip_ok": True}


def powerpoint_open_without_repair(path: str | Path) -> dict[str, Any]:
    """Windows PowerPoint COM oracle. Honest BLOCKED_EXTERNAL when COM is off."""
    if os.environ.get("ZECT_LIVE_PPT_COM", "").strip() != "1":
        return {"ok": False, "status": "BLOCKED_EXTERNAL", "reason": "ZECT_LIVE_PPT_COM!=1"}
    if os.name != "nt":
        return {"ok": False, "status": "BLOCKED_EXTERNAL", "reason": "not_windows"}
    try:
        import win32com.client  # type: ignore
    except ImportError:
        return {"ok": False, "status": "BLOCKED_EXTERNAL", "reason": "win32com_missing"}
    app = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(str(Path(path).resolve()), WithWindow=False)
        count = int(pres.Slides.Count)
        pres.Close()
        return {"ok": True, "status": "opened", "slide_count": count, "repair": False}
    except Exception as exc:  # noqa: BLE001 — COM oracle must not crash save
        return {"ok": False, "status": "repair_or_failed", "reason": str(exc)[:200]}
    finally:
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass


def _text_from_spec(spec: dict[str, Any]) -> str:
    parts: list[str] = []
    for block in spec.get("blocks") or []:
        if not isinstance(block, dict):
            continue
        kind = str(block.get("kind") or "")
        content = block.get("content") if isinstance(block.get("content"), dict) else {}
        if kind in {"text", "quote", "body", "title", "subtitle"}:
            blob = str(content.get("text") or "").strip()
            if blob:
                parts.append(blob)
    return "\n".join(parts) if parts else str(spec.get("text") or "").strip()


def _named_shape(slide, name: str):
    if not name:
        return None
    for shape in slide.shapes:
        if str(getattr(shape, "name", "") or "") == name:
            return shape
    return None


def _apply_named_block(slide, block: dict[str, Any]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    if content.get("locked"):
        return False
    name = str(content.get("shape_name") or "")
    shape = _named_shape(slide, name)
    if shape is None:
        return False
    geo = block.get("geometry") if isinstance(block.get("geometry"), dict) else None
    if geo:
        try:
            cx, cy = int(geo.get("cx") or 0), int(geo.get("cy") or 0)
            if cx > 0 and cy > 0:
                shape.left = Emu(int(geo.get("x") or 0))
                shape.top = Emu(int(geo.get("y") or 0))
                shape.width = Emu(cx)
                shape.height = Emu(cy)
        except Exception:
            pass
    kind = str(block.get("kind") or "")
    text = str(content.get("text") or "").strip()
    if kind in {"text", "quote", "body", "title", "subtitle"} and text and getattr(shape, "has_text_frame", False):
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        tf = shape.text_frame
        tf.clear()
        tf.text = (lines[0] if lines else text)[:800]
        for extra in lines[1:8]:
            p = tf.add_paragraph()
            p.text = extra[:800]
    return True


def _drop_pictures(slide) -> None:
    doomed = []
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                doomed.append(shape)
        except (ValueError, AttributeError):
            continue
    for shape in doomed:
        el = shape._element  # noqa: SLF001 — python-pptx has no public picture delete
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


def _update_table(slide, block: dict[str, Any]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    headers = [str(h) for h in list(content.get("headers") or [])]
    rows = [list(r) for r in list(content.get("rows") or []) if isinstance(r, list)]
    if not headers or not rows:
        return False
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        n_rows = min(len(table.rows), 1 + len(rows))
        n_cols = min(len(table.columns), len(headers))
        for j in range(n_cols):
            table.cell(0, j).text = headers[j][:80]
        for i in range(1, n_rows):
            row = rows[i - 1]
            for j in range(n_cols):
                table.cell(i, j).text = str(row[j] if j < len(row) else "")[:120]
        return True
    return False


def _apply_visual_blocks(slide, spec: dict[str, Any], *, user_id: str) -> None:
    from app.services.mentrix.presentation.charts import add_chart, replace_chart_data
    from app.services.mentrix.presentation.visual import paint_image, paint_table

    blocks = spec.get("blocks") if isinstance(spec.get("blocks"), list) else []
    default_geom = {"x": int(Inches(0.6)), "y": int(Inches(1.6)), "cx": int(Inches(8.8)), "cy": int(Inches(4.6))}
    replaced_image = False
    for block in blocks:
        if not isinstance(block, dict):
            continue
        kind = str(block.get("kind") or "")
        content = block.get("content") if isinstance(block.get("content"), dict) else {}
        if content.get("locked"):
            continue
        geom = block.get("geometry") if isinstance(block.get("geometry"), dict) else default_geom
        geom = {
            "x": int(geom.get("x") or default_geom["x"]),
            "y": int(geom.get("y") or default_geom["y"]),
            "cx": int(geom.get("cx") or default_geom["cx"]),
            "cy": int(geom.get("cy") or default_geom["cy"]),
        }
        if kind == "chart":
            if not replace_chart_data(slide, block):
                series = content.get("series") if isinstance(content.get("series"), list) else []
                if series:
                    add_chart(slide, block, geom)
        elif kind == "table":
            if not _update_table(slide, block):
                paint_table(slide, block, geom)
        elif kind == "image":
            if not content.get("asset_id"):
                continue
            if not replaced_image:
                _drop_pictures(slide)
                replaced_image = True
            paint_image(slide, block, geom, user_id=user_id)
        elif kind in {"shape", "metric", "quote", "diagram", "icon"}:
            from app.services.mentrix.presentation.visual import paint_block

            paint_block(slide, block, geom, user_id=user_id)


def _write_slide_text(slide, text: str) -> None:
    """Update existing title/body in place. Never add a covering dump textbox."""
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    title = lines[0][:160] if lines else text[:160]
    body = lines[1:]
    title_written = False
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        title_written = True
        rest = body
    else:
        rest = lines or [text]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            ph_type = shape.placeholder_format.type
        except ValueError:
            continue
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.VERTICAL_BODY):
            tf = shape.text_frame
            tf.clear()
            if rest:
                tf.text = rest[0][:800]
                for extra in rest[1:8]:
                    p = tf.add_paragraph()
                    p.text = extra[:800]
            return
    if not rest:
        return
    # Prefer rewriting an existing non-placeholder text frame rather than stacking a new box.
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            shape.placeholder_format.type
            continue
        except ValueError:
            current = (shape.text_frame.text or "").strip()
            if not current:
                continue
            if title_written and current[:80] == title[:80]:
                continue
            tf = shape.text_frame
            tf.clear()
            tf.text = rest[0][:800]
            for extra in rest[1:8]:
                p = tf.add_paragraph()
                p.text = extra[:800]
            return


def duplicate_slide_in_pptx(path: str | Path, index: int) -> dict[str, Any]:
    """Clone an existing slide in-place (OOXML shape tree copy)."""
    from copy import deepcopy

    pptx = Path(path)
    prs = Presentation(str(pptx))
    idx = int(index)
    if idx < 0 or idx >= len(prs.slides):
        raise ValueError("slide_index_out_of_range")
    source = prs.slides[idx]
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        newel = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")
    # Move new slide directly after source (python-pptx appends at end).
    slide_ids = prs.slides._sldIdLst
    new_id = slide_ids[-1]
    slide_ids.remove(new_id)
    slide_ids.insert(idx + 1, new_id)
    tmp = pptx.with_name(f"{pptx.stem}.zect-dup{pptx.suffix}")
    try:
        prs.save(str(tmp))
        os.replace(tmp, pptx)
        try:
            from app.services.mentrix.presentation.slide_preview import invalidate_slide_previews

            invalidate_slide_previews(pptx)
        except Exception:
            pass
    finally:
        if tmp.exists():
            tmp.unlink(missing_ok=True)
    sync_sidecar_from_pptx(pptx)
    return {"ok": True, "path": str(pptx), "slide_count": len(prs.slides), "inserted_at": idx + 1}


def sync_sidecar_from_pptx(pptx: Path) -> None:
    """Refresh notes sidecar from current PPTX so editor slide count matches file."""
    from app.services.mentrix.presentation.blocks import ensure_unique_block_ids
    from app.services.mentrix.presentation.document import document_from_pptx_bytes
    from app.services.pptx_paths import notes_sidecar_for_pptx, write_notes_sidecar

    doc = document_from_pptx_bytes(pptx.read_bytes(), path=str(pptx))
    slides = []
    for row in list(doc.get("slides") or []):
        if not isinstance(row, dict):
            continue
        slide_index = int(row.get("index") or len(slides))
        blocks = ensure_unique_block_ids(list(row.get("blocks") or []), slide_index=slide_index)
        slides.append(
            {
                "index": slide_index,
                "text": str(row.get("text") or ""),
                "notes": str(row.get("notes") or ""),
                "blocks": blocks,
            }
        )
    sidecar = notes_sidecar_for_pptx(pptx)
    write_notes_sidecar(sidecar, json.dumps({"path": str(pptx), "slides": slides}, indent=2))
