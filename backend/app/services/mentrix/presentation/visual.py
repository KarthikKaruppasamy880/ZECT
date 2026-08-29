"""Paint canonical visual blocks onto a python-pptx slide. Renderer calls this; it does not live in renderer helpers."""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.mentrix.presentation.asset_resolver import UnsafeImageError, load_image
from app.services.mentrix.presentation.charts import add_chart


def _box(slide, geometry: dict[str, int], *, fill: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(int(geometry["x"])),
        Emu(int(geometry["y"])),
        Emu(int(geometry["cx"])),
        Emu(int(geometry["cy"])),
    )
    shape.line.fill.background()
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def _text(shape, text: str, *, size: int = 14, bold: bool = False, color: RGBColor | None = None, align=PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = (text or "")[:800]
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def paint_image(slide, block: dict[str, Any], geometry: dict[str, int], *, user_id: str) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    asset_id = str(content.get("asset_id") or "").strip()
    if not asset_id:
        return False
    try:
        asset = load_image(asset_id, user_id=user_id)
    except (UnsafeImageError, FileNotFoundError, OSError):
        errors = list((block.get("validation") or {}).get("errors") or [])
        errors.append("image_load_failed")
        block["validation"] = {"ok": False, "errors": errors}
        return False
    slide.shapes.add_picture(
        str(asset["path"]),
        Emu(int(geometry["x"])),
        Emu(int(geometry["y"])),
        width=Emu(int(geometry["cx"])),
        height=Emu(int(geometry["cy"])),
    )
    caption = str(content.get("caption") or "").strip()
    if caption:
        cap_h = int(Inches(0.28))
        cap = {
            "x": geometry["x"],
            "y": geometry["y"] + max(geometry["cy"] - cap_h, 0),
            "cx": geometry["cx"],
            "cy": cap_h,
        }
        shape = _box(slide, cap, fill=RGBColor(0x00, 0x62, 0x8B))
        _text(shape, caption[:200], size=10, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    return True


def paint_table(slide, block: dict[str, Any], geometry: dict[str, int]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    headers = [str(h) for h in list(content.get("headers") or [])]
    rows = [list(r) for r in list(content.get("rows") or []) if isinstance(r, list)]
    if not headers or not rows:
        return False
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Emu(int(geometry["x"])),
        Emu(int(geometry["y"])),
        Emu(int(geometry["cx"])),
        Emu(int(geometry["cy"])),
    )
    table = table_shape.table
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header[:80]
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x62, 0x8B)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(11)
    for i, row in enumerate(rows):
        for j in range(n_cols):
            table.cell(i + 1, j).text = str(row[j] if j < len(row) else "")[:120]
    return True


def paint_metric(slide, block: dict[str, Any], geometry: dict[str, int]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    label = str(content.get("label") or "Metric").strip()
    value = str(content.get("value") or "—").strip()
    unit = str(content.get("unit") or "")
    if value.lower() in {"n/a", "na", "none", "-", "—"}:
        return False
    shape = _box(slide, geometry, fill=RGBColor(0xF7, 0xF4, 0xEF))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = label[:80]
    r0.font.size = Pt(12)
    r0.font.color.rgb = RGBColor(0x00, 0x62, 0x8B)
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = f"{value} {unit}".strip()[:60]
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xFF, 0x75, 0x00)
    delta = str(content.get("delta") or "").strip()
    if delta:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = delta[:40]
        r2.font.size = Pt(11)
    return True


def paint_quote(slide, block: dict[str, Any], geometry: dict[str, int]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    shape = _box(slide, geometry, fill=RGBColor(0xFF, 0xF4, 0xEB))
    text = str(content.get("text") or "").strip() or "Key message"
    attr = str(content.get("attribution") or "").strip()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = f"“{text[:400]}”"
    r0.font.size = Pt(18)
    r0.font.italic = True
    if attr:
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = f"— {attr[:80]}"
        r1.font.size = Pt(12)
        r1.font.color.rgb = RGBColor(0x00, 0x62, 0x8B)
    return True


def paint_diagram(slide, block: dict[str, Any], geometry: dict[str, int]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    nodes = [str(n) for n in list(content.get("nodes") or []) if str(n).strip()][:6]
    if not nodes:
        return False
    dtype = str(content.get("diagram_type") or "flow").lower()
    if dtype in {"process", "sequence"} and len(nodes) >= 2:
        row_cy = max(int(geometry["cy"] / max(len(nodes), 1) - int(Inches(0.08))), 280000)
        for i, node in enumerate(nodes):
            box = {
                "x": geometry["x"] + int(Inches(0.4)),
                "y": geometry["y"] + i * (row_cy + int(Inches(0.1))),
                "cx": geometry["cx"] - int(Inches(0.8)),
                "cy": row_cy,
            }
            shape = _box(slide, box, fill=RGBColor(0x00, 0x62, 0x8B))
            _text(shape, node[:60], size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        return True
    if dtype == "architecture" and len(nodes) >= 3:
        top, bottom = nodes[: (len(nodes) + 1) // 2], nodes[(len(nodes) + 1) // 2 :]
        rows = [top, bottom]
        row_cy = max(geometry["cy"] // 2 - int(Inches(0.12)), 350000)
        for r, row in enumerate(rows):
            gap = max(20000, geometry["cx"] // max(len(row) * 8, 1))
            cell_cx = max(int(geometry["cx"] / max(len(row), 1) - gap), 350000)
            y = geometry["y"] + r * (row_cy + int(Inches(0.2)))
            for i, node in enumerate(row):
                box = {"x": geometry["x"] + i * (cell_cx + gap), "y": y, "cx": cell_cx, "cy": row_cy}
                shape = _box(slide, box, fill=RGBColor(0x00, 0x62, 0x8B) if r == 0 else RGBColor(0x44, 0x54, 0x6A))
                _text(shape, node[:60], size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        return True
    gap = max(20000, geometry["cx"] // max(len(nodes) * 8, 1))
    cell_cx = max(int(geometry["cx"] / len(nodes) - gap), 320000)
    y = geometry["y"] + geometry["cy"] // 4
    cy = max(geometry["cy"] // 2, 360000)
    for i, node in enumerate(nodes):
        box = {"x": geometry["x"] + i * (cell_cx + gap), "y": y, "cx": cell_cx, "cy": cy}
        shape = _box(slide, box, fill=RGBColor(0x00, 0x62, 0x8B))
        _text(shape, node[:60], size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1 and dtype in {"flow", "process", "sequence", "timeline"}:
            arrow_x = box["x"] + box["cx"]
            arrow = {
                "x": arrow_x,
                "y": y + cy // 2 - int(Inches(0.08)),
                "cx": max(gap, 80000),
                "cy": int(Inches(0.16)),
            }
            _box(slide, arrow, fill=RGBColor(0xFF, 0x75, 0x00))
    return True


_ELEMENT_SHAPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "oval": MSO_SHAPE.OVAL,
    "arrow": MSO_SHAPE.RIGHT_ARROW,
}


def paint_shape(slide, block: dict[str, Any], geometry: dict[str, int]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    key = str(content.get("shape") or "rect").lower()
    mso = _ELEMENT_SHAPES.get(key, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(
        mso,
        Emu(int(geometry["x"])),
        Emu(int(geometry["y"])),
        Emu(int(geometry["cx"])),
        Emu(int(geometry["cy"])),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x62, 0x8B)
    shape.line.fill.background()
    label = str(content.get("text") or content.get("label") or key)[:80]
    if shape.has_text_frame and label:
        _text(shape, label, size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    return True


def paint_chart_svg_png(slide, block: dict[str, Any], geometry: dict[str, int]) -> bool:
    """Fallback raster for radar/area/stacked when python-pptx chart types fail."""
    import hashlib
    import struct
    import zlib
    from pathlib import Path
    from tempfile import gettempdir

    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    categories = [str(c) for c in list(content.get("categories") or [])][:8]
    series = [s for s in list(content.get("series") or []) if isinstance(s, dict)]
    values = [float(v) for v in list((series[0] or {}).get("values") or [])][: len(categories)] if series else []
    if len(categories) < 2 or len(values) < 2:
        return False
    w, h = 320, 200
    mx = max(values) or 1.0
    pts = []
    for i, val in enumerate(values):
        x = int(20 + i * ((w - 40) / max(len(values) - 1, 1)))
        y = int(h - 20 - (val / mx) * (h - 40))
        pts.append((x, y))
    # Minimal RGB PNG (dark teal background + orange polyline approximation as filled bars).
    raw = bytearray()
    for y in range(h):
        raw.append(0)
        for x in range(w):
            on = any(abs(x - px) < 3 and y >= py for px, py in pts)
            if on:
                raw.extend(b"\xff\x75\x00")
            else:
                raw.extend(b"\x00\x62\x8b")
    def _chunk(tag: bytes, data: bytes) -> bytes:
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

    png = b"\x89PNG\r\n\x1a\n" + _chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
    png += _chunk(b"IDAT", zlib.compress(bytes(raw), 9))
    png += _chunk(b"IEND", b"")
    digest = hashlib.sha1(png[:32]).hexdigest()[:12]
    path = Path(gettempdir()) / f"zect-chart-{digest}.png"
    path.write_bytes(png)
    slide.shapes.add_picture(
        str(path),
        Emu(int(geometry["x"])),
        Emu(int(geometry["y"])),
        width=Emu(int(geometry["cx"])),
        height=Emu(int(geometry["cy"])),
    )
    return True


def paint_icon(slide, block: dict[str, Any], geometry: dict[str, int]) -> bool:
    content = block.get("content") if isinstance(block.get("content"), dict) else {}
    glyph = str(content.get("glyph") or "★")[:4]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(geometry["x"])),
        Emu(int(geometry["y"])),
        Emu(int(geometry["cx"])),
        Emu(int(geometry["cy"])),
    )
    fill_hex = str(content.get("fill") or "#00628B").lstrip("#")
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(int(fill_hex[0:2], 16), int(fill_hex[2:4], 16), int(fill_hex[4:6], 16))
    except ValueError:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x00, 0x62, 0x8B)
    shape.line.fill.background()
    if shape.has_text_frame:
        _text(shape, glyph, size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    return True


def paint_block(slide, block: dict[str, Any], geometry: dict[str, int], *, user_id: str) -> bool:
    if not (block.get("validation") or {}).get("ok", True):
        errors = list((block.get("validation") or {}).get("errors") or [])
        if any(e in {"layout_overflow", "layout_collision", "layout_out_of_bounds", "image_url_rejected"} for e in errors):
            return False
        if "image_asset_required" in errors or "image_load_failed" in errors:
            return False
        if "chart_data_required" in errors or "chart_series_length_mismatch" in errors:
            return False
        if "table_data_required" in errors:
            return False
    kind = str(block.get("kind") or "")
    if kind == "image":
        return paint_image(slide, block, geometry, user_id=user_id)
    if kind == "chart":
        try:
            if add_chart(slide, block, geometry):
                return True
        except Exception:
            pass
        return paint_chart_svg_png(slide, block, geometry)
    if kind == "shape":
        return paint_shape(slide, block, geometry)
    if kind == "icon":
        return paint_icon(slide, block, geometry)
    if kind == "table":
        return paint_table(slide, block, geometry)
    if kind == "metric":
        return paint_metric(slide, block, geometry)
    if kind == "quote":
        return paint_quote(slide, block, geometry)
    if kind == "diagram":
        content = block.get("content") if isinstance(block.get("content"), dict) else {}
        if str(content.get("shape") or "") in _ELEMENT_SHAPES:
            return paint_shape(slide, block, geometry)
        return paint_diagram(slide, block, geometry)
    return False
