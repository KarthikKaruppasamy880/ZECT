"""Inspect the final PPTX/OOXML shape tree — not plan geometry."""

from __future__ import annotations

import io
import zipfile
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation

from app.services.mentrix.presentation.quality_policy import boxes_overlap, slide_size_emu

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _text_of(sp: ET.Element) -> str:
    parts = [t.text or "" for t in sp.findall(f".//{{{_NS_A}}}t")]
    return " ".join(p.strip() for p in parts if p and p.strip()).strip()


def _ph_type(sp: ET.Element) -> str | None:
    ph = sp.find(f".//{{{_NS_P}}}ph")
    if ph is None:
        return None
    return (ph.get("type") or "body").lower()


def _geom(sp: ET.Element) -> dict[str, int]:
    off = sp.find(f".//{{{_NS_A}}}off")
    ext = sp.find(f".//{{{_NS_A}}}ext")
    return {
        "x": int(off.get("x") or 0) if off is not None else 0,
        "y": int(off.get("y") or 0) if off is not None else 0,
        "cx": int(ext.get("cx") or 0) if ext is not None else 0,
        "cy": int(ext.get("cy") or 0) if ext is not None else 0,
    }


def _iter_shapes(slide_xml: bytes) -> list[dict[str, Any]]:
    root = ET.fromstring(slide_xml)
    out: list[dict[str, Any]] = []
    for sp in root.findall(f".//{{{_NS_P}}}sp"):
        nv = sp.find(f".//{{{_NS_P}}}cNvPr")
        name = nv.get("name") if nv is not None else ""
        geom = _geom(sp)
        text = _text_of(sp)
        out.append(
            {
                "name": name,
                "ph": _ph_type(sp),
                "geometry": geom,
                "text": text,
                "element": sp,
            }
        )
    return out


def _norm(text: str) -> str:
    return " ".join((text or "").lower().split())


def _semantic_duplicate_text(a: str, b: str) -> bool:
    """True when overlapping boxes likely duplicate content — not label + body pairs."""
    ta, tb = _norm(a), _norm(b)
    if not ta or not tb:
        return False
    if ta == tb:
        return True
    short, long = (ta, tb) if len(ta) <= len(tb) else (tb, ta)
    if short not in long:
        return False
    if len(short) < 20 and len(long) >= len(short) * 3:
        return False
    return len(short) / max(len(long), 1) >= 0.45


_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _broken_relationship_count(names: list[str], data: bytes) -> int:
    """Missing relationship targets (images/charts/media) are hard export blockers."""
    name_set = set(names)
    broken = 0
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            rels = [n for n in names if n.endswith(".rels")]
            for rel_path in rels:
                try:
                    root = ET.fromstring(zf.read(rel_path))
                except Exception:
                    broken += 1
                    continue
                parent = rel_path.rsplit("/_rels/", 1)[0] if "/_rels/" in rel_path else ""
                for rel in root.findall(f".//{{{_REL_NS}}}Relationship"):
                    target = (rel.get("Target") or "").replace("\\", "/")
                    if not target or target.startswith(("http://", "https://", "mailto:")):
                        continue
                    dest = target
                    if not dest.startswith("/"):
                        dest = f"{parent}/{target}" if parent else target
                    dest = dest.replace("/./", "/")
                    while "/../" in dest:
                        parts = dest.split("/")
                        acc: list[str] = []
                        for part in parts:
                            if part == "..":
                                if acc:
                                    acc.pop()
                            elif part != ".":
                                acc.append(part)
                        dest = "/".join(acc)
                    dest = dest.lstrip("/")
                    if dest not in name_set:
                        broken += 1
    except zipfile.BadZipFile:
        return 1
    return broken


def _resolved_slide_geoms(data: bytes) -> list[dict[str, dict[str, int]]]:
    """Placeholder xfrms often live on the layout; python-pptx resolves on-slide positions."""
    out: list[dict[str, dict[str, int]]] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = [i.filename.replace("\\", "/") for i in zf.infolist()]
        if not any(n.startswith("ppt/slideLayouts/") for n in names):
            return out
        prs = Presentation(io.BytesIO(data))
    except Exception:
        return out
    for slide in prs.slides:
        by_name: dict[str, dict[str, int]] = {}
        for shape in slide.shapes:
            try:
                name = str(getattr(shape, "name", "") or "")
                if not name:
                    continue
                by_name[name] = {
                    "x": int(shape.left or 0),
                    "y": int(shape.top or 0),
                    "cx": int(shape.width or 0),
                    "cy": int(shape.height or 0),
                }
            except Exception:
                continue
        out.append(by_name)
    return out


def _fill_missing_geom(shapes: list[dict[str, Any]], resolved: dict[str, dict[str, int]]) -> None:
    for shape in shapes:
        geom = shape.get("geometry") or {}
        if int(geom.get("cx") or 0) > 0 and int(geom.get("cy") or 0) > 0:
            continue
        hit = resolved.get(str(shape.get("name") or ""))
        if hit:
            shape["geometry"] = hit


def _slide_size_from_ooxml(data: bytes) -> tuple[int, int] | None:
    """Read p:sldSz from the package so 4:3 blanks are not judged against a 16:9 fallback."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            raw = zf.read("ppt/presentation.xml")
        root = ET.fromstring(raw)
    except Exception:
        return None
    for el in root.iter():
        if not str(el.tag).endswith("sldSz"):
            continue
        try:
            cx = int(el.get("cx") or 0)
            cy = int(el.get("cy") or 0)
        except (TypeError, ValueError):
            return None
        if cx > 0 and cy > 0:
            return cx, cy
    return None


def inspect_pptx_bytes(data: bytes, *, definition: dict[str, Any] | None = None) -> dict[str, Any]:
    if not data:
        return {
            "ok": False,
            "error": "empty_pptx",
            "status": "FAIL",
            "slides": [],
            "hard_findings": ["corrupt_pptx"],
            "export_blocked": True,
            "accept_warnings_allowed": False,
        }
    names: list[str] = []
    resolved_geoms = _resolved_slide_geoms(data)
    if definition is None:
        actual = _slide_size_from_ooxml(data)
        if actual:
            definition = {"slide_size": {"cx": actual[0], "cy": actual[1]}}
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = [i.filename.replace("\\", "/") for i in zf.infolist()]
            slide_names = sorted(
                n
                for n in names
                if n.startswith("ppt/slides/slide") and n.endswith(".xml") and "/_rels/" not in n
            )
            slides_out: list[dict[str, Any]] = []
            overlap_count = 0
            duplicate_title_count = 0
            dump_count = 0
            out_of_bounds = 0
            cx, cy = slide_size_emu(definition)
            for idx, name in enumerate(slide_names):
                shapes = _iter_shapes(zf.read(name))
                if idx < len(resolved_geoms):
                    _fill_missing_geom(shapes, resolved_geoms[idx])
                findings: list[str] = []
                text_boxes = [s for s in shapes if s["text"] or s["ph"]]
                titles = [s for s in text_boxes if (s["ph"] or "") in {"title", "ctrTitle", "title"} or "title" in (s["name"] or "").lower()]
                if len(titles) > 1:
                    texts = [_norm(s["text"]) for s in titles if s["text"]]
                    if len(texts) >= 2 and texts[0] and any(t == texts[0] or texts[0] in t or t in texts[0] for t in texts[1:]):
                        findings.append("duplicate_title")
                        duplicate_title_count += 1
                for i, a in enumerate(text_boxes):
                    ga = a["geometry"]
                    if ga["cx"] <= 0 or ga["cy"] <= 0:
                        continue
                    if ga["x"] + ga["cx"] > cx + 40000 or ga["y"] + ga["cy"] > cy + 40000:
                        out_of_bounds += 1
                        findings.append("out_of_bounds")
                    for b in text_boxes[i + 1 :]:
                        gb = b["geometry"]
                        if gb["cx"] <= 0 or gb["cy"] <= 0:
                            continue
                        if not boxes_overlap(ga, gb, pad=12000):
                            continue
                        ta, tb = _norm(a["text"]), _norm(b["text"])
                        if not ta or not tb:
                            continue
                        if _semantic_duplicate_text(ta, tb):
                            findings.append("duplicate_overlap")
                            overlap_count += 1
                dump = _find_dump_shape(shapes)
                if dump:
                    findings.append("covering_dump_textbox")
                    dump_count += 1
                ph_and_gen = _placeholder_plus_generated(shapes)
                if ph_and_gen:
                    findings.append("placeholder_and_generated")
                    overlap_count += 1
                slides_out.append(
                    {
                        "index": idx,
                        "shape_count": len(shapes),
                        "text_shape_count": len(text_boxes),
                        "findings": sorted(set(findings)),
                        "dump_shape": dump["name"] if dump else None,
                    }
                )
    except zipfile.BadZipFile:
        return {
            "ok": False,
            "error": "invalid_zip",
            "status": "FAIL",
            "slides": [],
            "hard_findings": ["corrupt_pptx"],
            "export_blocked": True,
            "accept_warnings_allowed": False,
        }
    broken_rels = _broken_relationship_count(names, data)
    hard_findings: list[str] = []
    if overlap_count:
        hard_findings.append("text_shape_collision")
    if duplicate_title_count:
        hard_findings.append("duplicate_semantic_content")
    if dump_count:
        hard_findings.append("covering_dump")
    if out_of_bounds:
        hard_findings.append("clipping_out_of_bounds")
    if broken_rels:
        hard_findings.append("broken_assets_or_rels")
    if "ppt/presentation.xml" not in names:
        hard_findings.append("corrupt_pptx")
    status = "FAIL" if hard_findings else "PASS"
    return {
        "ok": status == "PASS",
        "status": status,
        "slides": slides_out,
        "overlap_count": overlap_count,
        "duplicate_title_count": duplicate_title_count,
        "covering_dump_count": dump_count,
        "out_of_bounds_count": out_of_bounds,
        "broken_rel_count": broken_rels,
        "hard_findings": hard_findings,
        "export_blocked": bool(hard_findings),
        "accept_warnings_allowed": False if hard_findings else True,
        "slide_count": len(slides_out),
        "has_notes": sum(1 for n in names if n.startswith("ppt/notesSlides/")),
    }


def _find_dump_shape(shapes: list[dict[str, Any]]) -> dict[str, Any] | None:
    """A generated covering textbox whose text concatenates other shapes (zect-deck.pptx TextBox 10)."""
    others = [s for s in shapes if s["text"]]
    if len(others) < 3:
        return None
    for cand in shapes:
        if cand.get("ph"):
            continue
        geom = cand["geometry"]
        text = _norm(cand["text"])
        if not text or geom.get("cx", 0) < 7000000 or geom.get("cy", 0) < 3000000:
            continue
        covered = 0
        for other in others:
            if other is cand:
                continue
            ot = _norm(other["text"])
            if not ot or ot not in text:
                continue
            if boxes_overlap(geom, other["geometry"], pad=0) or geom.get("cy", 0) >= 3_000_000:
                covered += 1
        if covered >= 2:
            return cand
    return None


def _placeholder_plus_generated(shapes: list[dict[str, Any]]) -> bool:
    titles = [s for s in shapes if (s.get("ph") or "") in {"title", "ctrTitle"} and s.get("text")]
    gen_titles = [
        s
        for s in shapes
        if not s.get("ph") and s.get("text") and "title" in (s.get("name") or "").lower()
    ]
    if titles and gen_titles:
        t0 = _norm(titles[0]["text"])
        if any(_norm(g["text"]) == t0 or t0 in _norm(g["text"]) for g in gen_titles):
            return True
    placeholders = [s for s in shapes if s.get("ph") and s.get("text")]
    generated = [s for s in shapes if not s.get("ph") and s.get("text")]
    for ph in placeholders:
        for gen in generated:
            if not boxes_overlap(ph["geometry"], gen["geometry"], pad=12000):
                continue
            return True
    return False


def _shape_text_meta(shape) -> dict[str, Any] | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    text = (shape.text_frame.text or "").strip()
    if not text:
        return None
    is_ph = False
    try:
        _ = shape.placeholder_format.type
        is_ph = True
    except (ValueError, AttributeError):
        is_ph = False
    try:
        geom = {
            "x": int(shape.left or 0),
            "y": int(shape.top or 0),
            "cx": int(shape.width or 0),
            "cy": int(shape.height or 0),
        }
    except Exception:
        return None
    if geom["cx"] <= 0 or geom["cy"] <= 0:
        return None
    name = str(getattr(shape, "name", "") or "")
    return {
        "shape": shape,
        # OOXML's own shape id (<p:cNvPr id="...">) -- stable across
        # separate `slide.shapes` traversals, unlike Python's id(shape):
        # python-pptx builds a fresh wrapper object on every access to
        # `.shapes`, so two traversals of the same slide never share
        # object identity even though they wrap the same XML element.
        "shape_id": int(getattr(shape, "shape_id", 0) or 0),
        "name": name,
        "is_ph": is_ph,
        "text": text,
        "norm": _norm(text),
        "geom": geom,
        "area": geom["cx"] * geom["cy"],
    }


def _duplicate_text(a: str, b: str) -> bool:
    return _semantic_duplicate_text(a, b)


def _pick_overlap_victim(a: dict[str, Any], b: dict[str, Any]) -> dict[str, Any]:
    """Prefer keeping placeholders; drop generated duplicates and smaller boxes."""
    if a["is_ph"] and not b["is_ph"]:
        return b
    if b["is_ph"] and not a["is_ph"]:
        return a
    a_gen = "textbox" in a["name"].lower() and not a["is_ph"]
    b_gen = "textbox" in b["name"].lower() and not b["is_ph"]
    if a_gen and not b_gen:
        return a
    if b_gen and not a_gen:
        return b
    return a if a["area"] <= b["area"] else b


def strip_duplicate_overlapping_textboxes(data: bytes) -> tuple[bytes, int]:
    """Remove generated/duplicate text boxes that overlap placeholders or each other."""
    prs = Presentation(io.BytesIO(data))
    removed = 0
    for slide in prs.slides:
        metas = [m for m in (_shape_text_meta(sh) for sh in slide.shapes) if m]
        doomed: set[int] = set()
        for i, a in enumerate(metas):
            if a["shape_id"] in doomed:
                continue
            for b in metas[i + 1 :]:
                if b["shape_id"] in doomed:
                    continue
                if not boxes_overlap(a["geom"], b["geom"], pad=12000):
                    continue
                dup = _duplicate_text(a["norm"], b["norm"])
                ph_gen = (a["is_ph"] ^ b["is_ph"]) and (a["norm"] or b["norm"])
                if not dup and not ph_gen:
                    continue
                victim = _pick_overlap_victim(a, b)
                doomed.add(victim["shape_id"])
        for shape in list(slide.shapes):
            if int(getattr(shape, "shape_id", 0) or 0) not in doomed:
                continue
            el = shape._element  # noqa: SLF001 — python-pptx has no public delete
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)
                removed += 1
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), removed


def strip_covering_dump_shapes(data: bytes) -> tuple[bytes, int]:
    """Remove covering dump textboxes from a PPTX. Returns (bytes, removed_count)."""
    prs = Presentation(io.BytesIO(data))
    removed = 0
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        slide_names = sorted(
            n
            for n in zf.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml") and "/_rels/" not in n.replace("\\", "/")
        )
        dumps_by_index: dict[int, str] = {}
        for i, name in enumerate(slide_names):
            dump = _find_dump_shape(_iter_shapes(zf.read(name)))
            if dump:
                dumps_by_index[i] = dump["name"]
    for i, slide in enumerate(prs.slides):
        want = dumps_by_index.get(i)
        if not want:
            continue
        doomed = []
        for shape in slide.shapes:
            same_name = want and getattr(shape, "name", "") == want
            same_geom = False
            try:
                same_geom = int(shape.width or 0) >= 7_000_000 and int(shape.height or 0) >= 3_000_000
            except Exception:
                same_geom = False
            if same_name or (want and same_geom and not getattr(shape, "is_placeholder", False)):
                doomed.append(shape)
        for shape in doomed:
            el = shape._element  # noqa: SLF001
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)
                removed += 1
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), removed


def _broader_overlap_signal(data: bytes) -> int:
    """The same document-critic + rendered-quality overlap signal
    deck_catalog.quality_gate_for_path() folds into its overlap_total
    (max of inspector/doc/rendered overlap counts) -- inspect_and_repair_
    pptx() used to only ever look at its own narrower inspect_pptx_bytes()
    overlap_count, so it could (and did, on the zect-deck.pptx fixture)
    report "status": "PASS", "overlap_count": 0 for bytes the real export
    gate still failed with 15 overlaps it never even looked at. Never
    raises -- a critic failure must not block repair from returning at
    all, it just means this signal contributes 0 rather than crashing."""
    try:
        from app.services.mentrix.presentation.document import document_from_pptx_bytes
        from app.services.mentrix.presentation.quality_critic import critique_document

        critic = critique_document(document_from_pptx_bytes(data))
        rendered = critic.get("rendered_quality") if isinstance(critic.get("rendered_quality"), dict) else {}
        return max(
            int(critic.get("document_overlap_count") or 0),
            int(rendered.get("rendered_overlap_count") or 0),
        )
    except Exception:  # noqa: BLE001
        return 0


def inspect_and_repair_pptx(data: bytes, *, definition: dict[str, Any] | None = None) -> tuple[bytes, dict[str, Any]]:
    report = inspect_pptx_bytes(data, definition=definition)
    dump_removed = 0
    overlap_removed = 0
    if report.get("covering_dump_count"):
        data, dump_removed = strip_covering_dump_shapes(data)
        report = inspect_pptx_bytes(data, definition=definition)
    for _ in range(4):
        broader = _broader_overlap_signal(data)
        if not (report.get("overlap_count") or report.get("hard_findings") or broader):
            break
        before = max(int(report.get("overlap_count") or 0), broader)
        data, n = strip_duplicate_overlapping_textboxes(data)
        overlap_removed += n
        report = inspect_pptx_bytes(data, definition=definition)
        after = max(int(report.get("overlap_count") or 0), _broader_overlap_signal(data))
        if n == 0 or after >= before:
            break
    if dump_removed:
        report["dump_shapes_removed"] = dump_removed
    if overlap_removed:
        report["duplicate_shapes_removed"] = overlap_removed
    # Report the SAME broader signal the real export gate
    # (quality_gate_for_path) will check, so a caller of this function
    # alone never sees a false "status": "PASS" for bytes that gate would
    # still block.
    final_broader = _broader_overlap_signal(data)
    if final_broader > int(report.get("overlap_count") or 0):
        report["overlap_count"] = final_broader
        report["status"] = "FAIL"
        if "rendered_overlap" not in (report.get("hard_findings") or []):
            report.setdefault("hard_findings", []).append("rendered_overlap")
    report["final_artifact_status"] = report.get("status")
    return data, report


def merge_inspector_into_quality(quality: dict[str, Any], inspector: dict[str, Any]) -> dict[str, Any]:
    out = dict(quality or {})
    out["inspector"] = inspector
    out["overlap_count"] = max(int(out.get("overlap_count") or 0), int(inspector.get("overlap_count") or 0))
    out["out_of_bounds_count"] = max(int(out.get("out_of_bounds_count") or 0), int(inspector.get("out_of_bounds_count") or 0))
    out["covering_dump_count"] = int(inspector.get("covering_dump_count") or 0)
    out["duplicate_title_count"] = int(inspector.get("duplicate_title_count") or 0)
    layout_hard = (
        int(inspector.get("overlap_count") or 0)
        or int(inspector.get("covering_dump_count") or 0)
        or int(inspector.get("duplicate_title_count") or 0)
        or int(inspector.get("out_of_bounds_count") or 0)
        or int(inspector.get("broken_rel_count") or 0)
        or bool(inspector.get("hard_findings"))
    )
    if inspector.get("status") == "FAIL" or layout_hard:
        out["status"] = "FAIL"
        out["final_quality_status"] = "FAIL"
    out["layout_hard_fail"] = bool(layout_hard)
    return out
