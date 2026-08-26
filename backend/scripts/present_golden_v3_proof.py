#!/usr/bin/env python3
"""Golden V3 layout proof — human Zinnia + explicit 3 + unified quality."""

from __future__ import annotations

import json
import os
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

GOLDEN_PROMPT = (
    "Difference between AI Agentic and the Graph, loop and KV catch with LLM fine tuning"
)
REQUESTED = 3
TEMPLATE = "zinnia-executive-v1"


def main() -> int:
    os.environ.setdefault("ZECT_PRESENTATION_PROVIDER", "zect_native")
    from app.services.mentrix.presentation import prepare_prompt_deck
    from app.services.mentrix.presentation.deck_catalog import quality_gate_for_path
    from app.services.mentrix.presentation.native_provider import ZectNativePresentationProvider
    from app.services.mentrix.presentation.provider import PresentationGenerateRequest
    from app.services.pptx_parse import parse_pptx_bytes
    from pptx import Presentation

    prep = prepare_prompt_deck(prompt=GOLDEN_PROMPT, audience_id="general", requested_slide_count=REQUESTED)
    provider = ZectNativePresentationProvider()
    out = provider.generate(
        PresentationGenerateRequest(
            content=prep["adapted_prompt"],
            n_slides=REQUESTED,
            template=TEMPLATE,
            ui_template_choice=TEMPLATE,
            filename="golden-v3-agentic-deck.pptx",
            user_id="golden-v3-proof",
            fast_basic=True,
        )
    )
    artifact_dir = ROOT / "artifacts" / "present-golden-v3"
    artifact_dir.mkdir(parents=True, exist_ok=True)
    report: dict = {
        "ok": bool(out.get("ok")),
        "requested_slide_count": REQUESTED,
        "outline_target_line": "Target: 3 slides" in str(prep.get("adapted_prompt") or ""),
        "outline_no_audience_six": "Target ~6" not in str(prep.get("adapted_prompt") or ""),
        "plan_n_slides": out.get("n_slides"),
        "generation_job_id": out.get("generation_job_id"),
        "slide_count_trace": out.get("slide_count_trace"),
        "final_quality_status": out.get("final_quality_status"),
        "export_blocked": out.get("export_blocked"),
        "repair_attempts": out.get("repair_attempts"),
        "path": out.get("path"),
    }
    gate: dict = {}
    if out.get("path"):
        pptx = Path(str(out["path"]))
        data = pptx.read_bytes()
        parsed = parse_pptx_bytes(data)
        prs = Presentation(str(pptx))
        report["pptx_slide_count"] = len(prs.slides)
        report["parse_slide_count"] = len(parsed)
        report["slide_count_match"] = len(prs.slides) == REQUESTED
        gate = quality_gate_for_path(str(pptx))
        report["quality_gate"] = {
            "final_quality_status": gate.get("final_quality_status"),
            "rendered_overlap_count": gate.get("rendered_overlap_count"),
            "rendered_clipped_count": gate.get("rendered_clipped_count"),
            "template_conflict_count": gate.get("template_conflict_count"),
            "overlap_count": gate.get("overlap_count"),
            "export_blocked": gate.get("export_blocked"),
            "hard_findings": gate.get("hard_findings"),
        }
    report["acceptance"] = bool(
        out.get("ok")
        and report.get("slide_count_match")
        and int(out.get("n_slides") or 0) == REQUESTED
        and report.get("outline_target_line")
        and report.get("outline_no_audience_six")
        and str(gate.get("final_quality_status") or out.get("final_quality_status")) == "PASS"
        and int(gate.get("rendered_overlap_count") or 0) == 0
    )
    out_path = artifact_dir / "golden_v3_report.json"
    out_path.write_text(json.dumps(report, indent=2), encoding="utf-8")
    print(json.dumps(report, indent=2))
    return 0 if report["acceptance"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
